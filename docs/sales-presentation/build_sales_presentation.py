#!/usr/bin/env python3
"""
Roots säljpresentation för säljare som möter föreningar.
Nordisk minimalism, Roots brandbook. Fokus på "Why?".

Innehåll bygger på ASM-deckens säljargument + säljprocess (kollegans erfarenhet).
Bilder från apps/web/public (lätta att byta ut). Grundar-bilder är platshållare
(monogram-avatarer) – byt till riktiga foton när de finns.

Kör:  .venv/bin/python docs/sales-presentation/build_sales_presentation.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))
WEB = os.path.abspath(os.path.join(ROOT, "..", "..", "apps", "web", "public"))
SALES_ASSETS = os.path.abspath(os.path.join(ROOT, "..", "sales-material", "assets"))
CACHE = os.path.join(ROOT, "_cache")
os.makedirs(CACHE, exist_ok=True)
OUT = os.path.join(ROOT, "Roots_Saljpresentation.pptx")

# ───────────────────────── Brand ─────────────────────────
INK = RGBColor(0x1D, 0x1D, 0x1B)
SAND_DARK = RGBColor(0x7F, 0x71, 0x5B)
SAND_MED = RGBColor(0xB2, 0xA4, 0x91)
SAND_LIGHT = RGBColor(0xD5, 0xCA, 0xBF)
SAND_100 = RGBColor(0xF1, 0xEB, 0xE2)
OFFWHITE = RGBColor(0xFA, 0xF6, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOREST = RGBColor(0x6B, 0x79, 0x4F)
FOREST_SOFT = RGBColor(0xED, 0xF1, 0xE9)
OLIVE = RGBColor(0xC1, 0xBF, 0x99)
SKY = RGBColor(0xA7, 0xBB, 0xC5)
TERRA = RGBColor(0xE1, 0x87, 0x54)

HEAD = "Alan Sans"
BODY = "Inter 18pt"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_pageno = {"n": 0}


# ───────────────────────── Image cache (downscale to keep file small) ──────
def brand(name):
    """Return original PNG path (keeps transparency \u2013 don't JPEG-flatten logos)."""
    p = os.path.join(WEB, name)
    return p if os.path.exists(p) else None


def web_small(name, maxpx=1700, quality=82):
    """Downscale a web/public image into _cache and return abs path."""
    src = os.path.join(WEB, name)
    if not os.path.exists(src):
        return None
    safe = name.replace("/", "__")
    dst = os.path.join(CACHE, safe)
    if not os.path.exists(dst) or os.path.getmtime(dst) < os.path.getmtime(src):
        im = Image.open(src).convert("RGB")
        im.thumbnail((maxpx, maxpx))
        im.save(dst, "JPEG", quality=quality)
    return dst


# ───────────────────────── Primitives ─────────────────────────
def bg(slide, color=OFFWHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def rrect(slide, x, y, w, h, color, line=None, radius=0.06):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def oval(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1.5)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            text, font, size, color, bold = (seg + (False,))[:5] if len(seg) < 5 else seg
            r = p.add_run(); r.text = text
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold
    return tb


def kicker(slide, x, y, text, color=SAND_DARK, w=Inches(9)):
    return txt(slide, x, y, w, Inches(0.4),
               [[(text.upper(), BODY, 12, color, True)]], space_after=0)


def dot(slide, x, y, d, color):
    return oval(slide, x, y, d, d, color)


def img_fit(slide, abspath, x, y, max_w, max_h, align="center", valign="middle"):
    if not abspath or not os.path.exists(abspath):
        return None
    iw, ih = Image.open(abspath).size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = int(max_w / ar)
    else:
        h = max_h; w = int(max_h * ar)
    px = x + (max_w - w) // 2 if align == "center" else (x if align == "left" else x + (max_w - w))
    py = y + (max_h - h) // 2 if valign == "middle" else (y if valign == "top" else y + (max_h - h))
    return slide.shapes.add_picture(abspath, px, py, width=w, height=h)


def img_cover(slide, abspath, x, y, w, h):
    """Fill a box completely (center-crop), like CSS background-size: cover."""
    if not abspath or not os.path.exists(abspath):
        return None
    iw, ih = Image.open(abspath).size
    ar = iw / ih
    box_ar = w / h
    if ar > box_ar:
        nh = h; nw = int(h * ar)
    else:
        nw = w; nh = int(w / ar)
    pic = slide.shapes.add_picture(abspath, x - (nw - w) // 2, y - (nh - h) // 2,
                                   width=nw, height=nh)
    # crop overflow so it stays inside the box
    crop_lr = (nw - w) / nw / 2
    crop_tb = (nh - h) / nh / 2
    pic.crop_left = crop_lr; pic.crop_right = crop_lr
    pic.crop_top = crop_tb; pic.crop_bottom = crop_tb
    pic.left = x; pic.top = y; pic.width = w; pic.height = h
    return pic


def page_no(slide):
    _pageno["n"] += 1
    txt(slide, SW - Inches(0.95), SH - Inches(0.55), Inches(0.6), Inches(0.35),
        [[(str(_pageno["n"]), BODY, 10, SAND_MED, False)]], align=PP_ALIGN.RIGHT, space_after=0)


def footer(slide, light=False):
    c1 = OLIVE if light else SAND_DARK
    c2 = SAND_LIGHT if light else SAND_MED
    txt(slide, Inches(0.7), SH - Inches(0.55), Inches(9), Inches(0.35),
        [[("ROOTS", HEAD, 11, c1, True),
          ("   \u00b7   Säljpresentation \u2013 för föreningar", BODY, 10, c2, False)]],
        space_after=0)


def new_slide(color=OFFWHITE):
    s = prs.slides.add_slide(BLANK)
    bg(s, color)
    return s


def spine(slide, color=FOREST):
    rect(slide, 0, 0, Inches(0.28), SH, color)


def section_head(slide, kick, title, title_color=INK, y=Inches(0.7),
                 tw=Inches(12.0), tsize=30, kcolor=SAND_DARK):
    kicker(slide, Inches(0.7), y, kick, color=kcolor)
    txt(slide, Inches(0.66), y + Inches(0.42), tw, Inches(1.5),
        [[(title, HEAD, tsize, title_color, True)]], space_after=0, line_spacing=1.05)


def bullet_list(slide, items, x, y, w, gap=Inches(0.72), dotc=FOREST,
                tsize=14, bsize=12.5, tcolor=INK, bcolor=SAND_DARK):
    """items: list of (title, body or None)."""
    for i, it in enumerate(items):
        t, b = it if isinstance(it, tuple) else (it, None)
        cy = y + i * gap
        dot(slide, x, cy + Inches(0.07), Inches(0.16), dotc)
        if b:
            txt(slide, x + Inches(0.34), cy - Inches(0.04), w - Inches(0.34), Inches(0.4),
                [[(t, HEAD, tsize, tcolor, True)]], space_after=0)
            txt(slide, x + Inches(0.34), cy + Inches(0.28), w - Inches(0.34), Inches(0.5),
                [[(b, BODY, bsize, bcolor, False)]], line_spacing=1.25)
        else:
            txt(slide, x + Inches(0.34), cy - Inches(0.02), w - Inches(0.34), Inches(0.45),
                [[(t, BODY, tsize, tcolor, False)]], space_after=0, line_spacing=1.2)


# ════════════════════════ SLIDES ════════════════════════

# ── 1. Cover ────────────────────────────────────────────
def s_cover():
    s = new_slide(OFFWHITE)
    # höger: bild som fyller höjden
    img_cover(s, web_small("images/h1desktop.jpg", 2000), Inches(7.5), 0,
              SW - Inches(7.5), SH)
    rect(s, Inches(7.5), 0, Inches(0.06), SH, FOREST)
    spine(s, FOREST)
    img_fit(s, brand("brand/roots-logo-black.png"),
            Inches(0.9), Inches(0.85), Inches(2.4), Inches(0.9), align="left")
    txt(s, Inches(0.9), Inches(2.55), Inches(6.4), Inches(2.8),
        [[("Fyll föreningskassan", HEAD, 38, INK, True)],
         [("med produkter folk", HEAD, 38, INK, True)],
         [("faktiskt vill ha.", HEAD, 38, FOREST, True)]],
        space_after=2, line_spacing=1.05)
    txt(s, Inches(0.92), Inches(5.15), Inches(6.2), Inches(1.2),
        [[("Nordisk premium hår- & hudvård \u2013 bättre för håret, "
           "och byggd för att stärka svenskt föreningsliv.", BODY, 14.5, SAND_DARK, False)]],
        line_spacing=1.3)
    txt(s, Inches(0.92), Inches(6.55), Inches(6), Inches(0.5),
        [[("Säljpresentation \u00b7 möte med förening", BODY, 11.5, SAND_MED, False)]], space_after=0)


# ── 2. Vilka vi är (grundare – platshållare) ────────────
def founder(s, x, y, w, mono, name, role, bio, color):
    d = Inches(1.5)
    cx = x + (w - d) // 2
    oval(s, cx, y, d, d, color)
    txt(s, cx, y + Inches(0.42), d, Inches(0.7),
        [[(mono, HEAD, 30, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        space_after=0)
    txt(s, x, y + d + Inches(0.18), w, Inches(0.4),
        [[(name, HEAD, 16, INK, True)]], align=PP_ALIGN.CENTER, space_after=0)
    txt(s, x, y + d + Inches(0.55), w, Inches(0.35),
        [[(role.upper(), BODY, 10.5, FOREST, True)]], align=PP_ALIGN.CENTER, space_after=0)
    txt(s, x, y + d + Inches(0.92), w, Inches(1.3),
        [[(bio, BODY, 11.5, SAND_DARK, False)]], align=PP_ALIGN.CENTER, line_spacing=1.28)


def s_founders():
    s = new_slide(OFFWHITE)
    section_head(s, "Vilka vi är", "Vi kommer från föreningslivet")
    txt(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.7),
        [[("Vi har själva stått i klubbstugor, cuper, kiosker och säljmöten i många år. "
           "Vi vet vad som funkar \u2013 och vad som inte gör det.", BODY, 14, SAND_DARK, False)]],
        line_spacing=1.3)
    fy = Inches(3.05)
    fw = Inches(3.7); gap = Inches(0.35); x0 = Inches(0.9)
    founder(s, x0, fy, fw, "F1", "[Grundare 1]", "Föreningsförsäljning",
            "20+ år i föreningslivet och stora säljnätverk (Folkspel / Uppesittarkvällen). "
            "Vet hur en kassa fylls.", FOREST)
    founder(s, x0 + fw + gap, fy, fw, "F2", "[Grundare 2]", "Ledarskap & lagidrott",
            "Tränare och ledare. Förstår vad som driver lag, föräldrar och engagemang "
            "hela vägen ut.", SAND_DARK)
    founder(s, x0 + 2 * (fw + gap), fy, fw, "F3", "[Grundare 3]", "Produkt & digitalt",
            "Produktutveckling och den digitala plattformen \u2013 från formel till app "
            "och AI-stöd.", SKY)
    txt(s, Inches(0.7), SH - Inches(0.95), Inches(11.9), Inches(0.4),
        [[("Platshållarbilder \u2013 byts ut mot riktiga foton.", BODY, 10, SAND_MED, False)]],
        align=PP_ALIGN.CENTER, space_after=0)
    page_no(s)


# ── 3. Problem ──────────────────────────────────────────
def s_problem():
    s = new_slide(OFFWHITE)
    section_head(s, "Utgångsläget", "Föreningslivet förtjänar bättre")
    cards = [
        ("Mycket jobb", "Kataloger, blanketter och påminnelser. Stor insats för liten kassa."),
        ("Låg förtjänst", "Billiga produkter ger lite per sälj \u2013 det krävs väldigt många ordrar."),
        ("Krånglig admin", "Pengar, listor och leveranser landar hos redan trötta föräldrar."),
        ("Svårsålda produkter", "Strumpor, kakor och rabatthäften \u2013 sånt ingen längtar efter."),
    ]
    x = Inches(0.7); w = Inches(2.92); gap = Inches(0.18); y = Inches(2.75); h = Inches(3.05)
    for i, (t, b) in enumerate(cards):
        cx = x + i * (w + gap)
        rrect(s, cx, y, w, h, WHITE, line=SAND_LIGHT)
        dot(s, cx + Inches(0.32), y + Inches(0.32), Inches(0.3), TERRA)
        txt(s, cx + Inches(0.32), y + Inches(0.82), w - Inches(0.6), Inches(0.55),
            [[(t, HEAD, 16, INK, True)]], space_after=0)
        txt(s, cx + Inches(0.32), y + Inches(1.42), w - Inches(0.6), Inches(1.4),
            [[(b, BODY, 12, SAND_DARK, False)]], line_spacing=1.3)
    txt(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.7),
        [[("Frågan är inte ", BODY, 13.5, SAND_DARK, False),
          ("om", BODY, 13.5, FOREST, True),
          (" man ska sälja \u2013 utan ", BODY, 13.5, SAND_DARK, False),
          ("vad", BODY, 13.5, FOREST, True),
          (".", BODY, 13.5, SAND_DARK, False)]], line_spacing=1.3)
    footer(s); page_no(s)


# ── 4. Därför skapade vi Roots ──────────────────────────
def s_why_roots():
    s = new_slide(INK)
    spine(s, FOREST)
    kicker(s, Inches(0.7), Inches(0.75), "Vårt svar", color=OLIVE)
    txt(s, Inches(0.66), Inches(1.2), Inches(11.9), Inches(1.3),
        [[("Därför skapade vi Roots", HEAD, 32, WHITE, True)]], space_after=0)
    pts = [
        ("Premiumprodukter", "Hår- & hudvård folk faktiskt vill använda \u2013 varje dag."),
        ("35% tillbaka", "En tydlig, stark andel rakt in i föreningskassan."),
        ("Personligt stöd", "Vi håller er i handen hela vägen \u2013 före, under och efter."),
        ("Digital plattform", "Live-statistik, mål och betalning \u2013 utan pärmar."),
        ("Förenklad logistik", "Packat per förening, lag och spelare. Klart att dela ut."),
    ]
    x = Inches(0.7); y = Inches(2.85); w = Inches(5.7); colgap = Inches(0.6); rowgap = Inches(1.35)
    for i, (t, b) in enumerate(pts):
        cx = x + (i % 2) * (w + colgap)
        cy = y + (i // 2) * rowgap
        dot(s, cx, cy + Inches(0.06), Inches(0.26), OLIVE)
        txt(s, cx + Inches(0.42), cy - Inches(0.05), w - Inches(0.42), Inches(0.45),
            [[(t, HEAD, 17, WHITE, True)]], space_after=0)
        txt(s, cx + Inches(0.42), cy + Inches(0.38), w - Inches(0.42), Inches(0.85),
            [[(b, BODY, 12.5, SAND_LIGHT, False)]], line_spacing=1.28)
    page_no(s)


# ── 5. Why: borttagen marknadsföringskostnad ────────────
def s_why_marketing():
    s = new_slide(OFFWHITE)
    section_head(s, "Det här är vårt varför", "Vi tog bort det dyraste i beautybranschen")
    txt(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(0.8),
        [[("I traditionell hår- & hudvård går ", BODY, 14, SAND_DARK, False),
          ("20\u201340% av priset till marknadsföring", BODY, 14, INK, True),
          (" \u2013 annonser, influencers och kampanjer. Kostnader som inte gör produkten bättre.",
           BODY, 14, SAND_DARK, False)]], line_spacing=1.3)
    # två staplar
    bx = Inches(0.9); bw = Inches(5.4); bh = Inches(0.95)
    # Traditionellt
    y1 = Inches(3.35)
    txt(s, bx, y1 - Inches(0.45), bw, Inches(0.4),
        [[("Traditionellt varumärke", HEAD, 14, INK, True)]], space_after=0)
    rrect(s, bx, y1, bw, bh, SAND_100, radius=0.12)
    rrect(s, bx, y1, bw * 0.32, bh, TERRA, radius=0.12)
    txt(s, bx + Inches(0.2), y1, bw, bh,
        [[("20\u201340% marknadsföring", BODY, 12.5, WHITE, True)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    # Roots
    y2 = Inches(5.1)
    txt(s, bx, y2 - Inches(0.45), bw, Inches(0.4),
        [[("Roots", HEAD, 14, FOREST, True)]], space_after=0)
    rrect(s, bx, y2, bw, bh, FOREST_SOFT, radius=0.12)
    rrect(s, bx, y2, bw * 0.99, bh, FOREST, radius=0.12)
    txt(s, bx + Inches(0.2), y2, bw, bh,
        [[("Allt går till produkt + förening", BODY, 12.5, WHITE, True)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    # höger: slutsats
    rx = Inches(7.0)
    rrect(s, rx, Inches(3.0), Inches(5.6), Inches(3.55), INK)
    txt(s, rx + Inches(0.5), Inches(3.4), Inches(4.7), Inches(0.8),
        [[("Vi köper inte uppmärksamhet.", HEAD, 19, WHITE, True)]], line_spacing=1.1)
    txt(s, rx + Inches(0.5), Inches(4.25), Inches(4.7), Inches(0.8),
        [[("Vi bygger affärer genom engagemang.", HEAD, 19, OLIVE, True)]], line_spacing=1.1)
    txt(s, rx + Inches(0.5), Inches(5.35), Inches(4.7), Inches(1.0),
        [[("Pengarna går inte till reklam \u2013 de går till produkten och till er.",
           BODY, 13, SAND_LIGHT, False)]], line_spacing=1.3)
    footer(s); page_no(s)


# ── 6. En annan väg till kunden ─────────────────────────
def s_distribution():
    s = new_slide(OFFWHITE)
    section_head(s, "En annan väg till kunden", "Vi växer genom relationer \u2013 inte annonser")
    img_cover(s, web_small("images/02.jpg", 1800), Inches(7.2), Inches(2.1),
              Inches(5.45), Inches(4.5))
    pts = [
        ("Föreningarna driver försäljningen", "Lag, klasser och medlemmar \u2013 inte ett mediabolag."),
        ("Sälj via förtroende", "Man köper av någon man känner och vill stötta."),
        ("Direkt koppling", "Rak linje mellan kund, produkt och förmånstagare."),
    ]
    bullet_list(s, pts, Inches(0.7), Inches(2.55), Inches(6.0), gap=Inches(1.25),
                tsize=16, bsize=12.5)
    footer(s); page_no(s)


# ── 7. Mer värde där det räknas ─────────────────────────
def s_value_flow():
    s = new_slide(OFFWHITE)
    section_head(s, "Mer värde \u2013 där det faktiskt räknas",
                 "När vi slipper reklamkostnaden kan vi istället…")
    cards = [
        ("Bättre produkt", "Investera i kvalitet och upplevelse istället för räckvidd."),
        ("Attraktivt pris", "Premium folk vill köpa \u2013 till ett pris som funkar för familjer."),
        ("Starkare överskott", "Mer rakt in i föreningskassan, i varje såld produkt."),
        ("Långsiktig relation", "Återköp och kundrelationer som lever vidare över tid."),
    ]
    x = Inches(0.7); w = Inches(2.92); gap = Inches(0.18); y = Inches(2.65); h = Inches(2.95)
    cols = [FOREST, SAND_DARK, FOREST, SKY]
    for i, (t, b) in enumerate(cards):
        cx = x + i * (w + gap)
        rrect(s, cx, y, w, h, WHITE, line=SAND_LIGHT)
        rect(s, cx, y, w, Inches(0.12), cols[i])
        txt(s, cx + Inches(0.32), y + Inches(0.45), w - Inches(0.6), Inches(0.55),
            [[(t, HEAD, 16, INK, True)]], space_after=0)
        txt(s, cx + Inches(0.32), y + Inches(1.1), w - Inches(0.6), Inches(1.6),
            [[(b, BODY, 12, SAND_DARK, False)]], line_spacing=1.3)
    txt(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.7),
        [[("Andra bolag investerar i räckvidd. ", BODY, 13.5, SAND_DARK, False),
          ("Vi investerar i relationer.", BODY, 13.5, FOREST, True)]], line_spacing=1.3)
    footer(s); page_no(s)


# ── 8. Premiumprodukter folk vill köpa ──────────────────
def s_products_daily():
    s = new_slide(OFFWHITE)
    section_head(s, "Produkten", "Premium folk vill köpa \u2013 varje dag")
    img_cover(s, web_small("images/p4.jpg", 1800), Inches(6.9), Inches(2.0),
              Inches(5.75), Inches(4.7))
    pts = [
        ("Används varje dag", "Shampoo, balsam och body wash \u2013 inte en pryl som glöms bort."),
        ("Premiumkvalitet", "Skandinavisk känsla, ren formel, snygg design."),
        ("Återköp", "När flaskan tar slut kommer kunden tillbaka."),
        ("Stolthet att sälja", "Produkter man gärna sätter sitt namn på."),
    ]
    bullet_list(s, pts, Inches(0.7), Inches(2.5), Inches(5.8), gap=Inches(0.98),
                tsize=15, bsize=12)
    txt(s, Inches(0.7), Inches(6.55), Inches(5.9), Inches(0.6),
        [[("\u201cInte ännu en föreningsförsäljning \u2013 produkter folk faktiskt vill "
           "fortsätta köpa.\u201d", BODY, 11.5, FOREST, True)]], line_spacing=1.25)
    footer(s); page_no(s)


# ── 9. Premium till bättre pris ─────────────────────────
def s_price():
    s = new_slide(OFFWHITE)
    section_head(s, "Pris", "Premiumkänsla \u2013 till ett bättre pris")
    img_cover(s, web_small("images/p1.jpg", 1800), 0, Inches(2.1), Inches(5.6), Inches(4.6))
    # priskort
    rx = Inches(6.2)
    rrect(s, rx, Inches(2.4), Inches(6.4), Inches(1.7), FOREST)
    txt(s, rx + Inches(0.5), Inches(2.65), Inches(5.4), Inches(0.5),
        [[("Roots premiumpaket", BODY, 13, OLIVE, True)]], space_after=0)
    txt(s, rx + Inches(0.5), Inches(3.05), Inches(5.4), Inches(0.9),
        [[("399 kr", HEAD, 40, WHITE, True)]], space_after=0)
    rrect(s, rx, Inches(4.35), Inches(6.4), Inches(1.3), WHITE, line=SAND_LIGHT)
    txt(s, rx + Inches(0.5), Inches(4.55), Inches(5.4), Inches(0.5),
        [[("Motsvarande i butik", BODY, 13, SAND_DARK, True)]], space_after=0)
    txt(s, rx + Inches(0.5), Inches(4.9), Inches(5.4), Inches(0.7),
        [[("ofta 500\u2013700 kr", HEAD, 28, INK, True)]], space_after=0)
    txt(s, rx, Inches(5.95), Inches(6.4), Inches(0.8),
        [[("Premiumprodukter till ett attraktivt pris \u2013 samtidigt som man "
           "stöttar föreningslivet.", BODY, 12.5, SAND_DARK, False)]], line_spacing=1.3)
    footer(s); page_no(s)


# ── 10. Bättre för hår & hud (lätt vetenskap + USP) ─────
def s_science():
    s = new_slide(OFFWHITE)
    section_head(s, "Bättre för hår & hud", "Vi jobbar med hudbotten \u2013 inte emot den")
    img_cover(s, web_small("images/m2.jpg", 1600), 0, Inches(2.0), Inches(4.7), Inches(4.7))
    txt(s, Inches(5.1), Inches(2.1), Inches(7.5), Inches(0.9),
        [[("Håret växer ur en levande hudbotten. Roots är gjort för att stötta dess "
           "naturliga balans \u2013 mikrobiom och hudbottens eget lugnande system.",
           BODY, 13.5, SAND_DARK, False)]], line_spacing=1.3)
    usp = [
        ("SyriCalm", "Lugnar och balanserar hudbotten \u2013 certifierad aktiv (NATRUE)."),
        ("Panthenol (B5)", "Fukt och återfuktning som gör håret mjukare och starkare."),
        ("Sulfatsnålt & rent", "Rengör skonsamt, biologiskt nedbrytbart, ärlig formel."),
    ]
    bullet_list(s, usp, Inches(5.1), Inches(3.35), Inches(7.3), gap=Inches(0.98),
                tsize=15, bsize=12.5)
    txt(s, Inches(5.1), Inches(6.45), Inches(7.3), Inches(0.5),
        [[("Få ingredienser \u2013 de som finns gör jobbet. Fullständig förteckning i bilagan.",
           BODY, 11, SAND_MED, False)]], line_spacing=1.2)
    footer(s); page_no(s)


# ── 11. Jämförelse ──────────────────────────────────────
def s_compare():
    s = new_slide(OFFWHITE)
    section_head(s, "Vad skiljer oss", "Traditionell lagförsäljning vs Roots")
    rows = [
        ("Engångsköp", "Återköpsmöjlighet"),
        ("Kataloger & blanketter", "Digital plattform"),
        ("Begränsat stöd", "Dedikerad säljare som coachar"),
        ("Låg förtjänst", "35% tillbaka"),
        ("Svårsålda produkter", "Premium hår- & hudvård"),
    ]
    x = Inches(0.7); y = Inches(2.35)
    cw = Inches(5.85); gap = Inches(0.25); rh = Inches(0.72); rgap = Inches(0.12)
    # rubriker
    txt(s, x + Inches(0.3), y - Inches(0.05), cw, Inches(0.4),
        [[("Traditionellt", HEAD, 14, SAND_DARK, True)]], space_after=0)
    txt(s, x + cw + gap + Inches(0.3), y - Inches(0.05), cw, Inches(0.4),
        [[("Roots", HEAD, 14, FOREST, True)]], space_after=0)
    y0 = y + Inches(0.45)
    for i, (a, b) in enumerate(rows):
        ry = y0 + i * (rh + rgap)
        rrect(s, x, ry, cw, rh, SAND_100)
        txt(s, x + Inches(0.3), ry, cw - Inches(0.5), rh,
            [[(a, BODY, 13, SAND_DARK, False)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        rrect(s, x + cw + gap, ry, cw, rh, FOREST_SOFT)
        dot(s, x + cw + gap + Inches(0.28), ry + rh / 2 - Inches(0.08), Inches(0.16), FOREST)
        txt(s, x + cw + gap + Inches(0.6), ry, cw - Inches(0.8), rh,
            [[(b, BODY, 13, INK, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    footer(s); page_no(s)


# ── 12. Räknesnurran ────────────────────────────────────
def s_calculator():
    s = new_slide(INK)
    spine(s, FOREST)
    kicker(s, Inches(0.7), Inches(0.75), "Vår killer på mötet", color=OLIVE)
    txt(s, Inches(0.66), Inches(1.2), Inches(7.4), Inches(1.3),
        [[("Räknesnurran \u2013 låt ", HEAD, 30, WHITE, True),
          ("dem", HEAD, 30, OLIVE, True),
          (" säga siffrorna", HEAD, 30, WHITE, True)]], line_spacing=1.05)
    pts = [
        ("Två ingångar", "Per medlem (x paket/person) eller per lag (x paket/lag)."),
        ("Visualisera live", "Fyll i framför dem \u2013 summan växer i realtid."),
        ("Sätt förväntan", "\u201cVad tror du är rimligt per spelare?\u201d"),
        ("35% är låst", "Tydligt och ärligt \u2013 föreningen ser exakt vad den får."),
    ]
    for i, (t, b) in enumerate(pts):
        cy = Inches(2.85) + i * Inches(1.0)
        dot(s, Inches(0.7), cy + Inches(0.06), Inches(0.24), OLIVE)
        txt(s, Inches(1.1), cy - Inches(0.05), Inches(6.0), Inches(0.45),
            [[(t, HEAD, 15.5, WHITE, True)]], space_after=0)
        txt(s, Inches(1.1), cy + Inches(0.36), Inches(6.0), Inches(0.6),
            [[(b, BODY, 12, SAND_LIGHT, False)]], line_spacing=1.25)
    img_fit(s, web_small("demo/kalkylator-poster.jpg", 1400), Inches(8.4), Inches(1.4),
            Inches(4.3), Inches(5.6), valign="top")
    txt(s, Inches(8.4), SH - Inches(0.62), Inches(4.3), Inches(0.4),
        [[("Film & live-verktyg: roots.se/sa-fungerar-det", BODY, 10, SAND_MED, False)]],
        align=PP_ALIGN.CENTER, space_after=0)
    page_no(s)


# ── 13. WOW ─────────────────────────────────────────────
def s_wow():
    s = new_slide(FOREST)
    txt(s, Inches(0.9), Inches(0.85), Inches(11), Inches(0.5),
        [[("ETT RÄKNEEXEMPEL", BODY, 13, OLIVE, True)]], space_after=0)
    txt(s, Inches(0.86), Inches(1.4), Inches(11.5), Inches(2.0),
        [[("293 000 kr", HEAD, 80, WHITE, True)]], space_after=0)
    txt(s, Inches(0.9), Inches(3.35), Inches(11), Inches(0.6),
        [[("\u2013 tillbaka till föreningen.", HEAD, 24, OLIVE, True)]], space_after=0)
    items = [("300", "medlemmar"), ("8", "paket / person"), ("349 kr", "per paket"),
             ("35%", "tillbaka")]
    x = Inches(0.9); w = Inches(2.85); gap = Inches(0.2); y = Inches(4.6); h = Inches(1.7)
    for i, (big, small) in enumerate(items):
        cx = x + i * (w + gap)
        rrect(s, cx, y, w, h, OFFWHITE, radius=0.08)
        txt(s, cx, y + Inches(0.3), w, Inches(0.8),
            [[(big, HEAD, 34, INK, True)]], align=PP_ALIGN.CENTER, space_after=0)
        txt(s, cx, y + Inches(1.05), w, Inches(0.5),
            [[(small, BODY, 12.5, SAND_DARK, False)]], align=PP_ALIGN.CENTER, space_after=0)
    txt(s, Inches(0.9), SH - Inches(0.62), Inches(11), Inches(0.4),
        [[("Sätt siffrorna tillsammans med föreningen \u2013 målet blir deras eget.",
           BODY, 11.5, SAND_LIGHT, False)]], space_after=0)


# ── 14. Så enkelt är det (flöden) ───────────────────────
def s_flows():
    s = new_slide(OFFWHITE)
    section_head(s, "Plattformen", "Så enkelt är det \u2013 hela vägen")
    flows = [
        ("demo/forening-poster.jpg", "Föreningen", "Sätter mål, säljperiod och struktur \u2013 på minuter."),
        ("demo/lag-poster.jpg", "Lagledaren", "Bjuder in sitt lag och följer allt live."),
        ("demo/seller-poster.jpg", "Medlemmen", "Personlig shop, QR & länk \u2013 delar och säljer direkt."),
    ]
    x = Inches(0.9); w = Inches(3.7); gap = Inches(0.35); y = Inches(2.15)
    for i, (img, t, b) in enumerate(flows):
        cx = x + i * (w + gap)
        img_fit(s, web_small(img, 1100), cx, y, w, Inches(3.0), valign="top")
        txt(s, cx, y + Inches(3.1), w, Inches(0.4),
            [[(t, HEAD, 16, INK, True)]], align=PP_ALIGN.CENTER, space_after=0)
        txt(s, cx, y + Inches(3.5), w, Inches(0.9),
            [[(b, BODY, 12, SAND_DARK, False)]], align=PP_ALIGN.CENTER, line_spacing=1.25)
    txt(s, Inches(0.7), SH - Inches(0.62), Inches(11.9), Inches(0.4),
        [[("Se filmerna live: roots.se/sa-fungerar-det", BODY, 11, FOREST, True)]],
        align=PP_ALIGN.CENTER, space_after=0)
    page_no(s)


# ── 15. Logistik ────────────────────────────────────────
def s_logistics():
    s = new_slide(OFFWHITE)
    section_head(s, "Logistik", "Lägg tiden på ungdomarna \u2013 inte på kartonger")
    steps = ["Fabrik", "Roots lager\nGöteborg", "Packas per\nförening", "Packas per\nlag",
             "Packas per\nspelare", "Snabb\nutdelning"]
    n = len(steps); x0 = Inches(0.8); y = Inches(3.3)
    bw = Inches(1.72); gap = Inches(0.25)
    for i, st in enumerate(steps):
        cx = x0 + i * (bw + gap)
        c = FOREST if i in (0, n - 1) else WHITE
        tc = WHITE if i in (0, n - 1) else INK
        rrect(s, cx, y, bw, Inches(1.4), c, line=None if i in (0, n - 1) else SAND_LIGHT)
        txt(s, cx, y, bw, Inches(1.4),
            [[(line, HEAD, 13, tc, True)] for line in st.split("\n")],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.05)
        if i < n - 1:
            txt(s, cx + bw - Inches(0.02), y + Inches(0.45), gap + Inches(0.1), Inches(0.5),
                [[("\u2192", HEAD, 18, SAND_MED, True)]], align=PP_ALIGN.CENTER, space_after=0)
    txt(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.8),
        [[("Vi packar färdigt per spelare. Föreningen delar bara ut \u2013 inga listor, "
           "inga kartonger att sortera, ingen kontanthantering.", BODY, 14, SAND_DARK, False)]],
        line_spacing=1.3)
    footer(s); page_no(s)


# ── 16. Ni får en partner ───────────────────────────────
def s_partner():
    s = new_slide(OFFWHITE)
    section_head(s, "Det här ingår", "Ni får inte bara produkter \u2013 ni får en partner")
    phases = [
        ("FÖRE", ["Planering & leadslista", "Mål per medlem och lag", "Uppstart och kommunikation"]),
        ("UNDER", ["Coachning och push", "Live-tracking i appen", "\u201cNi är 72% till mål\u201d"]),
        ("EFTER", ["Resultat & topplista", "Färdigt content att dela", "Boka redan nästa period"]),
    ]
    x = Inches(0.7); w = Inches(3.85); gap = Inches(0.18); y = Inches(2.5); h = Inches(3.4)
    cols = [SAND_DARK, FOREST, SKY]
    for i, (ph, items) in enumerate(phases):
        cx = x + i * (w + gap)
        rrect(s, cx, y, w, h, WHITE, line=SAND_LIGHT)
        rect(s, cx, y, w, Inches(0.55), cols[i])
        txt(s, cx, y + Inches(0.07), w, Inches(0.45),
            [[(ph, HEAD, 15, WHITE, True)]], align=PP_ALIGN.CENTER, space_after=0)
        for j, it in enumerate(items):
            iy = y + Inches(0.85) + j * Inches(0.72)
            dot(s, cx + Inches(0.35), iy + Inches(0.06), Inches(0.14), cols[i])
            txt(s, cx + Inches(0.6), iy - Inches(0.04), w - Inches(0.85), Inches(0.6),
                [[(it, BODY, 12.5, SAND_DARK, False)]], line_spacing=1.2, space_after=0)
    txt(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.6),
        [[("Er framgång är vår framgång.", HEAD, 16, FOREST, True)]],
        align=PP_ALIGN.CENTER, space_after=0)
    footer(s); page_no(s)


# ── 17. Roots-appen ─────────────────────────────────────
def s_app():
    s = new_slide(INK)
    spine(s, FOREST)
    kicker(s, Inches(0.7), Inches(0.75), "Vi digitaliserar föreningsförsäljningen", color=OLIVE)
    txt(s, Inches(0.66), Inches(1.2), Inches(11.9), Inches(1.2),
        [[("Roots-appen", HEAD, 32, WHITE, True)]], space_after=0)
    chips = ["Live-statistik", "Ranking", "Säljtävlingar", "Gamification", "Coachning", "AI-stöd"]
    x = Inches(0.7); y = Inches(2.7); w = Inches(3.85); gap = Inches(0.2); h = Inches(1.15)
    for i, c in enumerate(chips):
        cx = x + (i % 3) * (w + gap)
        cy = y + (i // 3) * (h + Inches(0.25))
        rrect(s, cx, cy, w, h, RGBColor(0x2A, 0x2A, 0x27))
        dot(s, cx + Inches(0.32), cy + h / 2 - Inches(0.1), Inches(0.2), OLIVE)
        txt(s, cx + Inches(0.7), cy, w - Inches(0.9), h,
            [[(c, HEAD, 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    txt(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.6),
        [[("Energi, stolthet och driv \u2013 inbyggt i plattformen.", BODY, 13.5, SAND_LIGHT, False)]],
        space_after=0)
    page_no(s)


# ── 18. Vårt varför (värderingar) ───────────────────────
def s_values():
    s = new_slide(INK)
    img_cover(s, web_small("images/m3.jpg", 1800), Inches(8.0), 0, SW - Inches(8.0), SH)
    rect(s, Inches(8.0), 0, Inches(0.06), SH, FOREST)
    spine(s, FOREST)
    kicker(s, Inches(0.7), Inches(1.1), "Vårt varför", color=OLIVE)
    txt(s, Inches(0.66), Inches(1.7), Inches(7.0), Inches(2.4),
        [[("Vi vill bevisa att man kan", HEAD, 28, WHITE, True)],
         [("göra gott och sälja", HEAD, 28, OLIVE, True)],
         [("riktigt bra saker.", HEAD, 28, OLIVE, True)]],
        space_after=2, line_spacing=1.08)
    txt(s, Inches(0.7), Inches(4.35), Inches(6.7), Inches(2.4),
        [[("Varje flaska ger pengar till en förening i Sverige \u2013 och en bättre dag "
           "för någons hår. Inga genvägar på kvaliteten, inga överdrifter i löftena. "
           "Bara nordisk omtanke, från hudbotten och ut.", BODY, 15, SAND_LIGHT, False)]],
        line_spacing=1.4)
    page_no(s)


# ── 19. CTA / nästa steg ────────────────────────────────
def s_cta():
    s = new_slide(OFFWHITE)
    rect(s, 0, 0, SW, Inches(0.28), FOREST)
    img_fit(s, brand("brand/roots-symbol-black.png"), Inches(6.17), Inches(0.95),
            Inches(1.0), Inches(1.0))
    txt(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.3),
        [[("Nu kör vi.", HEAD, 36, INK, True)]], align=PP_ALIGN.CENTER, space_after=0)
    txt(s, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.8),
        [[("Vi sätter igång på under fem minuter \u2013 ingen startavgift, ingen "
           "bindningstid. Ni betalar först när produkterna är sålda och levererade.",
           BODY, 14, SAND_DARK, False)]], align=PP_ALIGN.CENTER, line_spacing=1.35)
    steps = ["Ansvarig / superuser", "Säljperiod & mål", "Kickoff och igång"]
    x = Inches(1.85); w = Inches(3.0); gap = Inches(0.42); y = Inches(4.5)
    for i, st in enumerate(steps):
        cx = x + i * (w + gap)
        rrect(s, cx, y, w, Inches(1.15), WHITE, line=SAND_LIGHT)
        oval(s, cx + Inches(0.3), y + Inches(0.32), Inches(0.5), Inches(0.5), FOREST)
        txt(s, cx + Inches(0.3), y + Inches(0.32), Inches(0.5), Inches(0.5),
            [[(str(i + 1), HEAD, 16, WHITE, True)]], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        txt(s, cx + Inches(0.95), y, w - Inches(1.1), Inches(1.15),
            [[(st, HEAD, 14, INK, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    rrect(s, Inches(5.16), Inches(6.1), Inches(3.0), Inches(0.7), FOREST)
    txt(s, Inches(5.16), Inches(6.19), Inches(3.0), Inches(0.55),
        [[("Starta er försäljning", BODY, 14, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    txt(s, Inches(1.0), Inches(7.0), Inches(11.3), Inches(0.4),
        [[("roots.se   \u00b7   hej@roots.se", BODY, 12, SAND_DARK, False)]],
        align=PP_ALIGN.CENTER, space_after=0)


# ════════════════ BILAGA ════════════════
def s_appendix_divider():
    s = new_slide(INK)
    spine(s, FOREST)
    txt(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(0.5),
        [[("BILAGA", BODY, 13, OLIVE, True)]], space_after=0)
    txt(s, Inches(0.66), Inches(3.5), Inches(11.9), Inches(1.4),
        [[("För dig som säljare", HEAD, 36, WHITE, True)]], space_after=0)
    txt(s, Inches(0.7), Inches(4.7), Inches(10), Inches(0.6),
        [[("Säljprocess, invändningar, produktfakta och kontakt.", BODY, 14, SAND_LIGHT, False)]],
        space_after=0)


def s_process():
    s = new_slide(OFFWHITE)
    section_head(s, "Bilaga \u00b7 playbook", "Säljprocessen i 7 steg")
    steps = [
        ("1 \u00b7 Lead & prioritering", "A/B/C. Sätt en hypotes redan innan mötet."),
        ("2 \u00b7 Hitta \u201cmotorn\u201d", "Boka rätt person \u2013 den som får saker att hända."),
        ("3 \u00b7 Säljmötet", "Insikt + potential = ja. Räknesnurran live."),
        ("4 \u00b7 Stängning & setup", "Ansvarig, säljperiod, leverans, struktur, mål."),
        ("5 \u00b7 Genomförande", "Push, live-tracking, gamification och stolthet."),
        ("6 \u00b7 Avslut & förstärkning", "Resultat, topplista, content \u2013 boka nästa."),
        ("7 \u00b7 Leverans & logistik", "Packat per spelare. Moment of truth."),
    ]
    x = Inches(0.7); w = Inches(5.85); colgap = Inches(0.3); y = Inches(2.2)
    rh = Inches(0.62); rgap = Inches(0.12)
    for i, (t, b) in enumerate(steps):
        col = i // 4
        row = i % 4
        cx = x + col * (w + colgap)
        cy = y + row * (rh + rgap + Inches(0.32))
        txt(s, cx, cy, w, Inches(0.4),
            [[(t, HEAD, 14, FOREST, True)]], space_after=0)
        txt(s, cx, cy + Inches(0.33), w, Inches(0.5),
            [[(b, BODY, 12, SAND_DARK, False)]], line_spacing=1.2, space_after=0)
    footer(s); page_no(s)


def s_objections():
    s = new_slide(OFFWHITE)
    section_head(s, "Bilaga \u00b7 playbook", "Invändningshantering")
    qa = [
        ("\u201cVi har redan en leverantör.\u201d",
         "Toppen \u2013 då vet ni hur mycket jobb det är. Roots ger högre förtjänst, "
         "premiumprodukter och vi sköter logistik och admin."),
        ("\u201cVåra föräldrar är trötta på att sälja.\u201d",
         "Precis därför. Allt är digitalt, packat per spelare, ingen kontanthantering. "
         "Mindre jobb \u2013 mer i kassan."),
        ("\u201c399 kr känns mycket.\u201d",
         "Motsvarande premium kostar 500\u2013700 kr i butik. Det är produkter folk "
         "ändå köper \u2013 nu stöttar de föreningen samtidigt."),
        ("\u201cTänk om vi inte säljer så mycket?\u201d",
         "Ingen risk: ingen startavgift, ni betalar först efter leverans. Vi sätter "
         "ett rimligt mål tillsammans och coachar er dit."),
    ]
    x = Inches(0.7); w = Inches(5.85); colgap = Inches(0.3); y = Inches(2.2)
    for i, (q, a) in enumerate(qa):
        cx = x + (i % 2) * (w + colgap)
        cy = y + (i // 2) * Inches(2.3)
        rrect(s, cx, cy, w, Inches(2.05), WHITE, line=SAND_LIGHT)
        txt(s, cx + Inches(0.3), cy + Inches(0.22), w - Inches(0.6), Inches(0.6),
            [[(q, HEAD, 14, INK, True)]], line_spacing=1.1)
        txt(s, cx + Inches(0.3), cy + Inches(0.85), w - Inches(0.6), Inches(1.1),
            [[(a, BODY, 12, SAND_DARK, False)]], line_spacing=1.25)
    footer(s); page_no(s)


def s_product_overview():
    s = new_slide(OFFWHITE)
    section_head(s, "Bilaga \u00b7 produkt", "Sortimentet")
    img_cover(s, web_small("images/p3.jpg", 1600), Inches(7.0), Inches(2.1),
              Inches(5.65), Inches(4.6))
    prods = [
        ("Hair Shampoo", "P2026110 \u00b7 Final Syricalm", "Rengör mjukt, balanserar hudbotten."),
        ("Conditioner", "P2026106 \u00b7 Pure less perfume", "Fukt och mjukhet, mindre parfym."),
        ("Body Wash", "P2026109 \u00b7 Final Syricalm", "Skonsam helkroppstvätt."),
    ]
    bullet_list(s, [(f"{n}  \u2013  {sub}", d) for n, sub, d in prods],
                Inches(0.7), Inches(2.55), Inches(6.0), gap=Inches(1.15),
                tsize=15, bsize=12.5)
    txt(s, Inches(0.7), Inches(6.3), Inches(6.0), Inches(0.6),
        [[("Fullständiga ingrediensförteckningar på följande sidor.", BODY, 11.5, SAND_MED, False)]],
        line_spacing=1.2)
    footer(s); page_no(s)


def s_ingredients(title, code, perfume, inci):
    s = new_slide(OFFWHITE)
    section_head(s, "Bilaga \u00b7 ingredienser", title)
    txt(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(0.5),
        [[(code + "   \u00b7   ", BODY, 13, FOREST, True),
          ("Parfym " + perfume, BODY, 13, SAND_DARK, False)]], space_after=0)
    rrect(s, Inches(0.7), Inches(2.6), Inches(11.93), Inches(3.9), WHITE, line=SAND_LIGHT)
    txt(s, Inches(1.05), Inches(2.9), Inches(11.25), Inches(3.3),
        [[("INCI", BODY, 11, SAND_MED, True)],
         [(inci, BODY, 13.5, INK, False)]], line_spacing=1.45, space_after=8)
    footer(s); page_no(s)


def s_company():
    s = new_slide(INK)
    spine(s, FOREST)
    img_fit(s, brand("brand/roots-symbol-white.png"), Inches(0.7), Inches(0.8),
            Inches(0.9), Inches(0.9), align="left")
    txt(s, Inches(0.66), Inches(2.0), Inches(11.9), Inches(1.0),
        [[("Företagsuppgifter", HEAD, 30, WHITE, True)]], space_after=0)
    txt(s, Inches(0.7), Inches(3.2), Inches(7), Inches(2.5),
        [[("Cafrelin AB", HEAD, 18, OLIVE, True)],
         [("c/o Fredrik Lindqvist", BODY, 15, SAND_LIGHT, False)],
         [("Hallängsvägen 8", BODY, 15, SAND_LIGHT, False)],
         [("439 55 Åsa", BODY, 15, SAND_LIGHT, False)],
         [("Org.nr 559355-7126", BODY, 15, SAND_LIGHT, False)]],
        line_spacing=1.35, space_after=4)
    txt(s, Inches(0.7), Inches(6.3), Inches(11), Inches(0.5),
        [[("roots.se   \u00b7   hej@roots.se", BODY, 13, SAND_MED, False)]], space_after=0)


def build():
    s_cover()
    s_founders()
    s_problem()
    s_why_roots()
    s_why_marketing()
    s_distribution()
    s_value_flow()
    s_products_daily()
    s_price()
    s_science()
    s_compare()
    s_calculator()
    s_wow()
    s_flows()
    s_logistics()
    s_partner()
    s_app()
    s_values()
    s_cta()
    # Bilaga
    s_appendix_divider()
    s_process()
    s_objections()
    s_product_overview()
    s_ingredients(
        "Body Wash Final Syricalm", "P2026109", "0,4%",
        "Aqua, Cocamidopropyl Betaine, Sodium Lauroyl Sarcosinate, Sodium Chloride, "
        "Citric Acid, Sodium Benzoate, Panthenyl Hydroxypropyl Steardimonium Chloride, "
        "PEG-150 Pentaerythrityl Tetrastearate, Parfum, Potassium Sorbate, "
        "PPG-2 Hydroxyethyl Cocamide, Panthenol, Sodium Citrate, Phragmites Communis Extract, "
        "Poria Cocos Extract.")
    s_ingredients(
        "Hair Shampoo Final Syricalm", "P2026110", "0,2%",
        "Aqua, Coco-Glucoside, Cocamidopropyl Betaine, Disodium Lauryl Sulfosuccinate, "
        "Glycerin, Sodium Chloride, PEG-4 Rapeseedamide, Sodium Benzoate, Citric Acid, "
        "Potassium Sorbate, Parfum, Polyquaternium-10, Polyquaternium-7, Sodium Citrate, "
        "Phragmites Communis Extract, Poria Cocos Extract, "
        "Octadecyl Di-t-Butyl-4-Hydroxyhydrocinnamate, Sodium Hydroxide.")
    s_ingredients(
        "Pure Conditioner less perfume Syricalm", "P2026106", "0,2%",
        "Aqua, Cetearyl Alcohol, Caprylic/Capric Triglyceride, Distearoylethyl "
        "Hydroxyethylmonium Methosulfate, Stearamidopropyl Dimethylamine, Phenoxyethanol, "
        "Panthenol, Hydrolyzed Corn Starch, Beta Vulgaris Root Extract, Butylene Glycol, "
        "Parfum, Citric Acid, Benzoic Acid, Sodium Lauroyl Lactylate, Sodium Caproyl Lactylate, "
        "Dehydroacetic Acid, Lactic Acid, Ethylhexylglycerin, Sodium Citrate, "
        "Piper Nigrum Fruit Extract, Phragmites Communis Extract, Poria Cocos Extract, "
        "Sodium Benzoate, Pentaerythrityl Tetra-Di-T-Butyl Hydroxyhydrocinnamate, "
        "Inga Alba Bark Extract, Tocopherol.")
    s_company()
    prs.save(OUT)
    print("Saved", os.path.relpath(OUT, ROOT), "-", len(prs.slides._sldIdLst), "slides")


build()
