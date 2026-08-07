# -*- coding: utf-8 -*-
"""
Roots presentationsmallar — återanvändbart mall-bibliotek (.pptx).

Nordisk minimalism, Roots brandbook. Samma typsnitt som hemsidan:
Alan Sans (rubrik/accent) + Inter (brödtext). 16:9.

Använd så här i ett byggskript:

    from roots_templates import Deck
    d = Deck(footer_label="Säljkickoff 2026")
    d.cover(kicker="...", title1="...", title2="...", subtitle="...")
    d.title_bullets(kicker="...", title="...", bullets=[...])
    ...
    d.save("Min_presentation.pptx")

Varje metod ritar en färdig slide och kan återanvändas i alla framtida
Roots-presentationer.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))
WEB = os.path.abspath(os.path.join(ROOT, "..", "..", "apps", "web", "public"))
CACHE = os.path.join(ROOT, "_cache")
os.makedirs(CACHE, exist_ok=True)

# ───────────────────────── Brand ─────────────────────────
INK = RGBColor(0x1D, 0x1D, 0x1B)
SAND_DARK = RGBColor(0x7F, 0x71, 0x5B)
SAND_MED = RGBColor(0xB2, 0xA4, 0x91)
SAND_LIGHT = RGBColor(0xD5, 0xCA, 0xBF)
OFFWHITE = RGBColor(0xFA, 0xF6, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOREST = RGBColor(0x6B, 0x79, 0x4F)
FOREST_SOFT = RGBColor(0xED, 0xF1, 0xE9)
OLIVE = RGBColor(0xC1, 0xBF, 0x99)

# Samma familjer som hemsidan (Alan Sans + Inter).
HEAD = "Alan Sans"
BODY = "Inter 18pt"


class Deck:
    """Ett Roots-presentationsdäck med färdiga, återanvändbara mallar."""

    def __init__(self, footer_label="Roots", width_in=13.333, height_in=7.5):
        self.prs = Presentation()
        self.prs.slide_width = Inches(width_in)
        self.prs.slide_height = Inches(height_in)
        self.SW = self.prs.slide_width
        self.SH = self.prs.slide_height
        self.BLANK = self.prs.slide_layouts[6]
        self.footer_label = footer_label
        self._pageno = 0

    # ───────────── Image helpers ─────────────
    @staticmethod
    def brand(name):
        p = os.path.join(WEB, name)
        return p if os.path.exists(p) else None

    @staticmethod
    def web_small(name, maxpx=1700, quality=82):
        src = os.path.join(WEB, name)
        if not os.path.exists(src):
            return None
        safe = name.replace("/", "__")
        dst = os.path.join(CACHE, safe)
        if not os.path.exists(dst) or os.path.getmtime(dst) < os.path.getmtime(src):
            im = Image.open(src).convert("RGB")
            im.thumbnail((maxpx, maxpx))
            im.save(dst, "JPEG", quality=quality)
        return dst

    # ───────────── Primitives ─────────────
    def _bg(self, slide, color=OFFWHITE):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _rect(self, slide, x, y, w, h, color, line=None):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        if line:
            sp.line.color.rgb = line; sp.line.width = Pt(1)
        else:
            sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def _rrect(self, slide, x, y, w, h, color, line=None, radius=0.06):
        sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        if line:
            sp.line.color.rgb = line; sp.line.width = Pt(1)
        else:
            sp.line.fill.background()
        sp.shadow.inherit = False
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
        return sp

    def _oval(self, slide, x, y, w, h, color, line=None):
        sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        if line:
            sp.line.color.rgb = line; sp.line.width = Pt(1.5)
        else:
            sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def _txt(self, slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, space_after=6, line_spacing=1.15):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        if isinstance(runs[0], tuple):
            runs = [runs]
        for i, para in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(space_after)
            p.space_before = Pt(0)
            p.line_spacing = line_spacing
            for seg in para:
                text, font, size, color, bold = (seg + (False,))[:5] if len(seg) < 5 else seg
                r = p.add_run(); r.text = text
                r.font.name = font; r.font.size = Pt(size)
                r.font.color.rgb = color; r.font.bold = bold
        return tb

    def _kicker(self, slide, x, y, text, color=SAND_DARK, w=Inches(10)):
        return self._txt(slide, x, y, w, Inches(0.4),
                         [[(text.upper(), BODY, 12, color, True)]], space_after=0)

    def _dot(self, slide, x, y, d, color):
        return self._oval(slide, x, y, d, d, color)

    def _img_fit(self, slide, abspath, x, y, max_w, max_h, align="center", valign="middle"):
        if not abspath or not os.path.exists(abspath):
            return None
        iw, ih = Image.open(abspath).size
        ar = iw / ih
        box_ar = max_w / max_h
        if ar > box_ar:
            w = max_w; h = int(max_w / ar)
        else:
            h = max_h; w = int(max_h * ar)
        px = x + (max_w - w) // 2 if align == "center" else (x if align == "left" else x + (max_w - w))
        py = y + (max_h - h) // 2 if valign == "middle" else (y if valign == "top" else y + (max_h - h))
        return slide.shapes.add_picture(abspath, px, py, width=w, height=h)

    def _img_cover(self, slide, abspath, x, y, w, h):
        if not abspath or not os.path.exists(abspath):
            return None
        iw, ih = Image.open(abspath).size
        ar = iw / ih
        box_ar = w / h
        if ar > box_ar:
            nh = h; nw = int(h * ar)
        else:
            nw = w; nh = int(w / ar)
        pic = slide.shapes.add_picture(abspath, x - (nw - w) // 2, y - (nh - h) // 2,
                                       width=nw, height=nh)
        crop_lr = (nw - w) / nw / 2
        crop_tb = (nh - h) / nh / 2
        pic.crop_left = crop_lr; pic.crop_right = crop_lr
        pic.crop_top = crop_tb; pic.crop_bottom = crop_tb
        pic.left = x; pic.top = y; pic.width = w; pic.height = h
        return pic

    def _spine(self, slide, color=FOREST):
        self._rect(slide, 0, 0, Inches(0.28), self.SH, color)

    def _page_no(self, slide, light=False):
        self._pageno += 1
        c = SAND_LIGHT if light else SAND_MED
        self._txt(slide, self.SW - Inches(0.95), self.SH - Inches(0.55),
                  Inches(0.6), Inches(0.35),
                  [[(str(self._pageno), BODY, 10, c, False)]],
                  align=PP_ALIGN.RIGHT, space_after=0)

    def _footer(self, slide, light=False):
        c1 = OLIVE if light else SAND_DARK
        c2 = SAND_LIGHT if light else SAND_MED
        self._txt(slide, Inches(0.7), self.SH - Inches(0.55), Inches(9), Inches(0.35),
                  [[("ROOTS", HEAD, 11, c1, True),
                    ("   \u00b7   " + self.footer_label, BODY, 10, c2, False)]],
                  space_after=0)

    def _new(self, color=OFFWHITE):
        s = self.prs.slides.add_slide(self.BLANK)
        self._bg(s, color)
        return s

    def _section_head(self, slide, kick, title, y=Inches(0.72), tsize=32, tw=Inches(11.6)):
        self._kicker(slide, Inches(0.7), y, kick)
        self._txt(slide, Inches(0.66), y + Inches(0.42), tw, Inches(1.4),
                  [[(title, HEAD, tsize, INK, True)]], space_after=0, line_spacing=1.03)

    def _cards_geom(self, n, left=Inches(0.7), gutter=Inches(0.35)):
        right = self.SW - Inches(0.7)
        total = right - left
        cw = int((total - gutter * (n - 1)) / n)
        xs = [left + i * (cw + gutter) for i in range(n)]
        return xs, cw

    def tag(self, slide, text):
        """Liten etikett uppe till höger (används i mall-katalogen)."""
        w = Inches(3.2)
        self._txt(slide, self.SW - w - Inches(0.7), Inches(0.4), w, Inches(0.35),
                  [[(text.upper(), BODY, 10, SAND_MED, True)]],
                  align=PP_ALIGN.RIGHT, space_after=0)

    # ════════════════════ MALLAR / TEMPLATES ════════════════════
    def cover(self, kicker, title1, title2, subtitle, image="images/sport-hero.jpg"):
        """MALL: Titelsida — logo, rubrik i två delar, underrubrik, bild höger."""
        s = self._new(OFFWHITE)
        img = self.web_small(image, 2000) if image else None
        if img:
            self._img_cover(s, img, Inches(7.7), 0, self.SW - Inches(7.7), self.SH)
            self._rect(s, Inches(7.7), 0, Inches(0.06), self.SH, FOREST)
        self._spine(s, FOREST)
        self._img_fit(s, self.brand("brand/roots-logo-black.png"),
                      Inches(0.9), Inches(0.85), Inches(2.4), Inches(0.9), align="left")
        self._kicker(s, Inches(0.92), Inches(2.35), kicker)
        self._txt(s, Inches(0.9), Inches(2.75), Inches(6.6), Inches(2.4),
                  [[(title1, HEAD, 40, INK, True)],
                   [(title2, HEAD, 40, FOREST, True)]],
                  space_after=2, line_spacing=1.04)
        self._txt(s, Inches(0.92), Inches(4.85), Inches(6.2), Inches(1.4),
                  [[(subtitle, BODY, 15, SAND_DARK, False)]], line_spacing=1.35)
        return s

    def divider(self, num, title):
        """MALL: Avsnittsdelare — mörk, stort nummer + rubrik."""
        s = self._new(INK)
        self._spine(s, OLIVE)
        self._img_fit(s, self.brand("brand/roots-symbol-white.png"),
                      self.SW - Inches(1.7), Inches(0.7), Inches(0.9), Inches(0.9), align="right")
        self._txt(s, Inches(0.9), Inches(2.4), Inches(6), Inches(1.6),
                  [[(num, HEAD, 84, OLIVE, True)]], space_after=0, line_spacing=1.0)
        self._rect(s, Inches(0.98), Inches(3.95), Inches(1.4), Inches(0.05), FOREST)
        self._txt(s, Inches(0.95), Inches(4.15), Inches(10.5), Inches(1.6),
                  [[(title, HEAD, 34, WHITE, True)]], space_after=0, line_spacing=1.05)
        self._footer(s, light=True); self._page_no(s, light=True)
        return s

    def title_only(self, kicker, title, tsize=40):
        """MALL: Endast rubrik — stort påstående, mycket luft."""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._kicker(s, Inches(0.9), Inches(2.7), kicker, color=FOREST)
        self._txt(s, Inches(0.88), Inches(3.15), Inches(11.4), Inches(2.6),
                  [[(title, HEAD, tsize, INK, True)]], space_after=0, line_spacing=1.06)
        self._footer(s); self._page_no(s)
        return s

    def title_text(self, kicker, title, body):
        """MALL: Rubrik + textruta — rubrik och ett brödtextstycke."""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        body_runs = [[(p, BODY, 15, SAND_DARK, False)] for p in
                     (body if isinstance(body, list) else [body])]
        self._txt(s, Inches(0.7), Inches(2.7), Inches(8.6), Inches(4),
                  body_runs, space_after=10, line_spacing=1.4)
        self._footer(s); self._page_no(s)
        return s

    def title_bullets(self, kicker, title, bullets):
        """MALL: Rubrik + punktlista."""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        n = len(bullets)
        top = Inches(2.7) if n <= 3 else Inches(2.5)
        gap = Inches(1.0) if n <= 3 else Inches(0.82)
        for i, b in enumerate(bullets):
            cy = top + i * gap
            self._dot(s, Inches(0.74), cy + Inches(0.09), Inches(0.2), FOREST)
            self._txt(s, Inches(1.2), cy, Inches(10.8), Inches(0.8),
                      [[(b, HEAD, 19, INK, True)]], space_after=0, line_spacing=1.05)
        self._footer(s); self._page_no(s)
        return s

    def agenda(self, kicker, title, items):
        """MALL: Agenda — numrerad lista i två kolumner."""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        col_x = [Inches(0.72), Inches(6.9)]
        top = Inches(2.5); gap = Inches(0.92); per_col = 5
        for i, it in enumerate(items):
            col = 0 if i < per_col else 1
            row = i if i < per_col else i - per_col
            x = col_x[col]; cy = top + row * gap
            self._txt(s, x, cy - Inches(0.02), Inches(0.7), Inches(0.6),
                      [[(f"{i+1:02d}", HEAD, 20, OLIVE, True)]], space_after=0)
            self._txt(s, x + Inches(0.78), cy + Inches(0.03), Inches(5.2), Inches(0.7),
                      [[(it, HEAD, 16, INK, True)]], space_after=0, line_spacing=1.05)
        self._footer(s); self._page_no(s)
        return s

    def title_image(self, kicker, title, body, image, img_side="right"):
        """MALL: Rubrik + bild + textruta — text och bild sida vid sida."""
        s = self._new(OFFWHITE)
        self._spine(s)
        img = self.web_small(image, 1600)
        half = Inches(5.55)
        if img_side == "right":
            self._section_head(s, kicker, title)
            self._img_cover(s, img, self.SW - half - Inches(0.7), Inches(1.4),
                            half, Inches(5.1))
            self._txt(s, Inches(0.7), Inches(2.7), Inches(5.6), Inches(4),
                      [[(p, BODY, 14.5, SAND_DARK, False)] for p in
                       (body if isinstance(body, list) else [body])],
                      space_after=10, line_spacing=1.4)
        else:
            self._img_cover(s, img, Inches(0.7), Inches(0.7), half, self.SH - Inches(1.4))
            tx = Inches(6.6)
            self._kicker(s, tx, Inches(0.9), kicker)
            self._txt(s, tx - Inches(0.04), Inches(1.32), Inches(6), Inches(1.4),
                      [[(title, HEAD, 30, INK, True)]], space_after=0, line_spacing=1.05)
            self._txt(s, tx, Inches(2.9), Inches(6), Inches(4),
                      [[(p, BODY, 14.5, SAND_DARK, False)] for p in
                       (body if isinstance(body, list) else [body])],
                      space_after=10, line_spacing=1.4)
        self._footer(s); self._page_no(s)
        return s

    def image_full(self, title, subtitle, image):
        """MALL: Helbild med rubrik-overlay."""
        s = self._new(INK)
        img = self.web_small(image, 2200)
        self._img_cover(s, img, 0, 0, self.SW, self.SH)
        # mörk gradient-liknande platta i botten för läsbarhet
        self._rect(s, 0, self.SH - Inches(2.4), self.SW, Inches(2.4), INK)
        try:
            s.shapes[-1].fill.fore_color.rgb = INK
            s.shapes[-1].fill.transparency = 0  # platt – python-pptx saknar äkta alpha
        except Exception:
            pass
        self._txt(s, Inches(0.9), self.SH - Inches(1.95), Inches(11.5), Inches(1.1),
                  [[(title, HEAD, 34, WHITE, True)]], space_after=0, line_spacing=1.05)
        self._txt(s, Inches(0.92), self.SH - Inches(1.05), Inches(11), Inches(0.6),
                  [[(subtitle, BODY, 14, OLIVE, False)]], space_after=0)
        self._page_no(s, light=True)
        return s

    def two_col(self, kicker, title, left_title, left_body, right_title, right_body):
        """MALL: Två kolumner — jämför/ställ mot varandra."""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        cy = Inches(2.7); cw = Inches(5.75); ch = Inches(3.9)
        for x, ht, bd in [(Inches(0.7), left_title, left_body),
                          (Inches(6.85), right_title, right_body)]:
            self._rrect(s, x, cy, cw, ch, WHITE, radius=0.05)
            self._rect(s, x + Inches(0.34), cy + Inches(0.42), Inches(0.5), Inches(0.06), FOREST)
            self._txt(s, x + Inches(0.34), cy + Inches(0.66), cw - Inches(0.68), Inches(0.7),
                      [[(ht, HEAD, 20, INK, True)]], space_after=0, line_spacing=1.05)
            self._txt(s, x + Inches(0.34), cy + Inches(1.45), cw - Inches(0.68), Inches(2.2),
                      [[(p, BODY, 13.5, SAND_DARK, False)] for p in
                       (bd if isinstance(bd, list) else [bd])],
                      space_after=8, line_spacing=1.35)
        self._footer(s); self._page_no(s)
        return s

    def cards(self, kicker, title, cards):
        """MALL: Kort (2–4 st) — rubrik + kort med titel och text."""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        n = len(cards)
        xs, cw = self._cards_geom(n)
        cy = Inches(2.7); ch = Inches(3.4)
        for (t, b), x in zip(cards, xs):
            self._rrect(s, x, cy, cw, ch, WHITE, radius=0.05)
            self._rect(s, x + Inches(0.32), cy + Inches(0.42), Inches(0.5), Inches(0.06), FOREST)
            self._txt(s, x + Inches(0.32), cy + Inches(0.68), cw - Inches(0.64), Inches(0.9),
                      [[(t, HEAD, 19, INK, True)]], space_after=0, line_spacing=1.05)
            self._txt(s, x + Inches(0.32), cy + Inches(1.55), cw - Inches(0.64), Inches(1.6),
                      [[(b, BODY, 13, SAND_DARK, False)]], line_spacing=1.32)
        self._footer(s); self._page_no(s)
        return s

    def stat(self, kicker, title, stats):
        """MALL: Nyckeltal — stora siffror i rad. stats=[(tal, etikett), ...]"""
        s = self._new(FOREST_SOFT)
        self._spine(s, FOREST)
        self._kicker(s, Inches(0.9), Inches(0.9), kicker, color=FOREST)
        self._txt(s, Inches(0.86), Inches(1.32), Inches(11.4), Inches(1.0),
                  [[(title, HEAD, 30, INK, True)]], space_after=0, line_spacing=1.03)
        n = len(stats)
        xs, cw = self._cards_geom(n)
        cy = Inches(3.1)
        for (num, label), x in zip(stats, xs):
            self._txt(s, x, cy, cw, Inches(1.4),
                      [[(num, HEAD, 60, FOREST, True)]], align=PP_ALIGN.CENTER,
                      space_after=0, line_spacing=1.0)
            self._txt(s, x, cy + Inches(1.5), cw, Inches(0.9),
                      [[(label, BODY, 14, SAND_DARK, False)]], align=PP_ALIGN.CENTER,
                      space_after=0, line_spacing=1.25)
        self._footer(s); self._page_no(s)
        return s

    def hero(self, kicker, line1, line2):
        """MALL: Statement — tvåradig rubrik i två färger på mjuk bakgrund."""
        s = self._new(FOREST_SOFT)
        self._spine(s, FOREST)
        self._kicker(s, Inches(0.9), Inches(2.2), kicker, color=FOREST)
        self._txt(s, Inches(0.88), Inches(2.7), Inches(11.4), Inches(3.4),
                  [[(line1, HEAD, 40, INK, True)],
                   [(line2, HEAD, 40, FOREST, True)]],
                  space_after=4, line_spacing=1.08)
        self._footer(s); self._page_no(s)
        return s

    def quote(self, text, attribution=None):
        """MALL: Citat / statement — stor centrerad text på mjuk bakgrund."""
        s = self._new(FOREST_SOFT)
        self._spine(s, FOREST)
        self._txt(s, Inches(1.4), Inches(2.4), Inches(10.5), Inches(3.0),
                  [[("\u201c" + text + "\u201d", HEAD, 34, INK, True)]],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  space_after=0, line_spacing=1.12)
        if attribution:
            self._txt(s, Inches(1.4), Inches(5.4), Inches(10.5), Inches(0.6),
                      [[("— " + attribution, BODY, 14, SAND_DARK, False)]],
                      align=PP_ALIGN.CENTER, space_after=0)
        self._footer(s); self._page_no(s)
        return s

    def products(self, kicker, title, cards,
                 imgs=("images/sport-schampoo.jpg", "images/sport-conditioner.jpg", "images/sport-body-wash.jpg")):
        """MALL: Produktkort — bild + typ + namn + fördel. cards=[(namn, typ, fördel)]"""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        n = len(cards)
        xs, cw = self._cards_geom(n)
        cy = Inches(2.55); ch = Inches(3.7)
        for (name, typ, benefit), x, im in zip(cards, xs, imgs):
            self._rrect(s, x, cy, cw, ch, WHITE, radius=0.05)
            self._img_cover(s, self.web_small(im, 1200), x, cy, cw, Inches(1.75))
            self._txt(s, x + Inches(0.3), cy + Inches(1.95), cw - Inches(0.6), Inches(0.5),
                      [[(typ.upper(), BODY, 10.5, FOREST, True)]], space_after=0)
            self._txt(s, x + Inches(0.3), cy + Inches(2.3), cw - Inches(0.6), Inches(0.6),
                      [[(name, HEAD, 20, INK, True)]], space_after=0, line_spacing=1.0)
            self._txt(s, x + Inches(0.3), cy + Inches(2.95), cw - Inches(0.6), Inches(0.7),
                      [[(benefit, BODY, 12.5, SAND_DARK, False)]], line_spacing=1.3)
        self._footer(s); self._page_no(s)
        return s

    def flow(self, kicker, title, steps):
        """MALL: Flöde / steg — numrerade rutor med pilar."""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        left = Inches(0.72); right = self.SW - Inches(0.72)
        per_row = 4; gutter = Inches(0.3)
        total = right - left
        cw = int((total - gutter * (per_row - 1)) / per_row)
        ch = Inches(1.35)
        rows_y = [Inches(2.9), Inches(4.75)]
        for i, st in enumerate(steps):
            row = i // per_row; col = i % per_row
            x = left + col * (cw + gutter); y = rows_y[row]
            self._rrect(s, x, y, cw, ch, FOREST_SOFT, radius=0.12)
            self._txt(s, x + Inches(0.28), y + Inches(0.2), cw - Inches(0.5), Inches(0.4),
                      [[(f"{i+1:02d}", HEAD, 15, OLIVE, True)]], space_after=0)
            self._txt(s, x + Inches(0.28), y + Inches(0.62), cw - Inches(0.5), Inches(0.6),
                      [[(st, HEAD, 15, INK, True)]], space_after=0, line_spacing=1.0)
            if col < per_row - 1 and i + 1 < len(steps):
                ax = x + cw + Inches(0.04)
                self._txt(s, ax, y + Inches(0.42), gutter - Inches(0.02), Inches(0.5),
                          [[("\u203a", HEAD, 22, SAND_MED, True)]],
                          align=PP_ALIGN.CENTER, space_after=0)
        self._footer(s); self._page_no(s)
        return s

    def screens(self, kicker, title, screens):
        """MALL: Skärmar / mockups — bild i mörk ram + bildtext. screens=[(bild, rubrik, undertext)]"""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        n = len(screens)
        xs, cw = self._cards_geom(n)
        top = Inches(2.45); frame_h = Inches(3.4)
        for (img, cap, sub), x in zip(screens, xs):
            self._rrect(s, x, top, cw, frame_h, INK, radius=0.06)
            self._img_fit(s, self.web_small(img, 1200), x + Inches(0.12), top + Inches(0.12),
                          cw - Inches(0.24), frame_h - Inches(0.24), valign="top")
            self._txt(s, x, top + frame_h + Inches(0.18), cw, Inches(0.5),
                      [[(cap, HEAD, 16, INK, True)]], align=PP_ALIGN.CENTER, space_after=0)
            self._txt(s, x, top + frame_h + Inches(0.58), cw, Inches(0.5),
                      [[(sub, BODY, 12, SAND_DARK, False)]], align=PP_ALIGN.CENTER, space_after=0)
        self._footer(s); self._page_no(s)
        return s

    def closing(self, title, subtitle):
        """MALL: Avslutning — mörk, stor rubrik + underrubrik."""
        s = self._new(INK)
        self._spine(s, OLIVE)
        self._img_fit(s, self.brand("brand/roots-symbol-white.png"),
                      Inches(0.9), Inches(0.85), Inches(0.95), Inches(0.95), align="left")
        self._txt(s, Inches(0.88), Inches(2.9), Inches(11.4), Inches(1.4),
                  [[(title, HEAD, 54, WHITE, True)]], space_after=0, line_spacing=1.0)
        self._txt(s, Inches(0.92), Inches(4.35), Inches(11), Inches(0.8),
                  [[(subtitle, BODY, 17, OLIVE, False)]], space_after=0)
        return s

    def big_number(self, kicker, number, label, sub=None):
        """MALL: Rubrik + stor siffra — en enda siffra som får bära sidan."""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._kicker(s, Inches(0.9), Inches(1.1), kicker, color=FOREST)
        self._txt(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.6),
                  [[(number, HEAD, 150, FOREST, True)]], align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
        self._txt(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(0.8),
                  [[(label, HEAD, 26, INK, True)]], align=PP_ALIGN.CENTER,
                  space_after=0, line_spacing=1.0)
        if sub:
            self._txt(s, Inches(1.9), Inches(5.7), Inches(9.5), Inches(0.9),
                      [[(sub, BODY, 14, SAND_DARK, False)]], align=PP_ALIGN.CENTER,
                      space_after=0, line_spacing=1.35)
        self._footer(s); self._page_no(s)
        return s

    def timeline(self, kicker, title, milestones):
        """MALL: Tidslinje — horisontell linje med hållpunkter.
        milestones=[(etikett, text), ...] (3–5 st)"""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        n = len(milestones)
        left = Inches(1.85); right = self.SW - Inches(1.75)
        line_y = Inches(4.15)
        self._rect(s, left, line_y, right - left, Inches(0.035), SAND_MED)
        step = (right - left) / (n - 1) if n > 1 else 0
        for i, (label, text) in enumerate(milestones):
            cx = int(left + i * step)
            d = Inches(0.26)
            self._oval(s, cx - d // 2, line_y - d // 2 + Emu(17000), d, d, FOREST)
            colw = Inches(2.6)
            self._txt(s, cx - colw // 2, line_y - Inches(1.35), colw, Inches(0.6),
                      [[(label, HEAD, 17, FOREST, True)]], align=PP_ALIGN.CENTER,
                      space_after=0, line_spacing=1.0)
            self._txt(s, cx - colw // 2, line_y + Inches(0.4), colw, Inches(1.6),
                      [[(text, BODY, 12.5, SAND_DARK, False)]], align=PP_ALIGN.CENTER,
                      space_after=0, line_spacing=1.3)
        self._footer(s); self._page_no(s)
        return s

    def table(self, kicker, title, headers, rows):
        """MALL: Tabell — jämförelse/data. headers=[...], rows=[[...], ...]"""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        ncols = len(headers); nrows = len(rows) + 1
        left = Inches(0.7); top = Inches(2.65)
        width = self.SW - Inches(1.4)
        rowh = min(Inches(0.62), Inches(3.9) / nrows) if nrows else Inches(0.62)
        height = int(rowh) * nrows
        gf = s.shapes.add_table(nrows, ncols, left, top, width, height)
        tbl = gf.table
        tbl.first_row = True
        tbl.horz_banding = False
        for j, h in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = h
            c.fill.solid(); c.fill.fore_color.rgb = FOREST
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.18); c.margin_right = Inches(0.1)
            p = c.text_frame.paragraphs[0]
            r = p.runs[0] if p.runs else p.add_run()
            r.font.name = BODY; r.font.size = Pt(12.5)
            r.font.bold = True; r.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                c = tbl.cell(i, j)
                c.text = str(val)
                c.fill.solid()
                c.fill.fore_color.rgb = WHITE if i % 2 else FOREST_SOFT
                c.vertical_anchor = MSO_ANCHOR.MIDDLE
                c.margin_left = Inches(0.18); c.margin_right = Inches(0.1)
                p = c.text_frame.paragraphs[0]
                r = p.runs[0] if p.runs else p.add_run()
                r.font.name = BODY; r.font.size = Pt(12)
                r.font.bold = (j == 0)
                r.font.color.rgb = INK if j == 0 else SAND_DARK
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        self._footer(s); self._page_no(s)
        return s

    def team(self, kicker, title, members):
        """MALL: Team — foton med namn + roll. members=[(bild, namn, roll)]
        Byt ut bilderna mot riktiga porträtt när de finns."""
        s = self._new(OFFWHITE)
        self._spine(s)
        self._section_head(s, kicker, title)
        n = len(members)
        xs, cw = self._cards_geom(n)
        top = Inches(2.7)
        photo_h = Inches(3.0)
        for (img, name, role), x in zip(members, xs):
            self._rrect(s, x, top, cw, photo_h, INK, radius=0.06)
            self._img_cover(s, self.web_small(img, 1400),
                            x + Emu(1), top + Emu(1), cw - Emu(2), photo_h - Emu(2))
            self._txt(s, x, top + photo_h + Inches(0.22), cw, Inches(0.5),
                      [[(name, HEAD, 19, INK, True)]], align=PP_ALIGN.CENTER, space_after=0)
            self._txt(s, x, top + photo_h + Inches(0.66), cw, Inches(0.5),
                      [[(role.upper(), BODY, 11, FOREST, True)]],
                      align=PP_ALIGN.CENTER, space_after=0)
        self._footer(s); self._page_no(s)
        return s

    # ───────────── Save ─────────────
    def save(self, path, embed=True):
        self.prs.save(path)
        if embed:
            try:
                from pptx_embed import embed_fonts
                embed_fonts(path)
            except Exception as e:  # embedding är en bonus, aldrig blockerande
                print("  (font-embedding hoppades över:", e, ")")
        return path
