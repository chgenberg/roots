#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Roots säljkickoff – intern keynote (.pptx) + talarmanus/stödord (.docx).

Nordisk minimalism, Roots brandbook, låg textmängd på slides – detaljerna
ligger i stödorden. Innehållet kommer från content.py (en sanning).

Kör:  .venv/bin/python docs/sales-kickoff/build_kickoff.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

from docx import Document
from docx.shared import Pt as DPt, RGBColor as DColor, Inches as DInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from content import SLIDES

ROOT = os.path.dirname(os.path.abspath(__file__))
WEB = os.path.abspath(os.path.join(ROOT, "..", "..", "apps", "web", "public"))
CACHE = os.path.join(ROOT, "_cache")
os.makedirs(CACHE, exist_ok=True)
OUT_PPTX = os.path.join(ROOT, "Roots_Saljkickoff.pptx")
OUT_DOCX = os.path.join(ROOT, "Roots_Saljkickoff_Talarmanus.docx")

# ───────────────────────── Brand ─────────────────────────
INK = RGBColor(0x1D, 0x1D, 0x1B)
SAND_DARK = RGBColor(0x7F, 0x71, 0x5B)
SAND_MED = RGBColor(0xB2, 0xA4, 0x91)
SAND_LIGHT = RGBColor(0xD5, 0xCA, 0xBF)
OFFWHITE = RGBColor(0xFA, 0xF6, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOREST = RGBColor(0x6B, 0x79, 0x4F)
FOREST_SOFT = RGBColor(0xED, 0xF1, 0xE9)
OLIVE = RGBColor(0xC1, 0xBF, 0x99)

HEAD = "Alan Sans"
BODY = "Inter 18pt"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_pageno = {"n": 0}


# ───────────────────────── Image helpers ─────────────────────────
def brand(name):
    p = os.path.join(WEB, name)
    return p if os.path.exists(p) else None


def web_small(name, maxpx=1700, quality=82):
    src = os.path.join(WEB, name)
    if not os.path.exists(src):
        return None
    safe = name.replace("/", "__")
    dst = os.path.join(CACHE, safe)
    if not os.path.exists(dst) or os.path.getmtime(dst) < os.path.getmtime(src):
        im = Image.open(src).convert("RGB")
        im.thumbnail((maxpx, maxpx))
        im.save(dst, "JPEG", quality=quality)
    return dst


# ───────────────────────── Primitives ─────────────────────────
def bg(slide, color=OFFWHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def rrect(slide, x, y, w, h, color, line=None, radius=0.06):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def oval(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1.5)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            text, font, size, color, bold = (seg + (False,))[:5] if len(seg) < 5 else seg
            r = p.add_run(); r.text = text
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold
    return tb


def kicker(slide, x, y, text, color=SAND_DARK, w=Inches(10)):
    return txt(slide, x, y, w, Inches(0.4),
               [[(text.upper(), BODY, 12, color, True)]], space_after=0)


def dot(slide, x, y, d, color):
    return oval(slide, x, y, d, d, color)


def img_fit(slide, abspath, x, y, max_w, max_h, align="center", valign="middle"):
    if not abspath or not os.path.exists(abspath):
        return None
    iw, ih = Image.open(abspath).size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = int(max_w / ar)
    else:
        h = max_h; w = int(max_h * ar)
    px = x + (max_w - w) // 2 if align == "center" else (x if align == "left" else x + (max_w - w))
    py = y + (max_h - h) // 2 if valign == "middle" else (y if valign == "top" else y + (max_h - h))
    return slide.shapes.add_picture(abspath, px, py, width=w, height=h)


def img_cover(slide, abspath, x, y, w, h):
    if not abspath or not os.path.exists(abspath):
        return None
    iw, ih = Image.open(abspath).size
    ar = iw / ih
    box_ar = w / h
    if ar > box_ar:
        nh = h; nw = int(h * ar)
    else:
        nw = w; nh = int(w / ar)
    pic = slide.shapes.add_picture(abspath, x - (nw - w) // 2, y - (nh - h) // 2,
                                   width=nw, height=nh)
    crop_lr = (nw - w) / nw / 2
    crop_tb = (nh - h) / nh / 2
    pic.crop_left = crop_lr; pic.crop_right = crop_lr
    pic.crop_top = crop_tb; pic.crop_bottom = crop_tb
    pic.left = x; pic.top = y; pic.width = w; pic.height = h
    return pic


def spine(slide, color=FOREST):
    rect(slide, 0, 0, Inches(0.28), SH, color)


def page_no(slide, light=False):
    _pageno["n"] += 1
    c = SAND_LIGHT if light else SAND_MED
    txt(slide, SW - Inches(0.95), SH - Inches(0.55), Inches(0.6), Inches(0.35),
        [[(str(_pageno["n"]), BODY, 10, c, False)]], align=PP_ALIGN.RIGHT, space_after=0)


def footer(slide, light=False):
    c1 = OLIVE if light else SAND_DARK
    c2 = SAND_LIGHT if light else SAND_MED
    txt(slide, Inches(0.7), SH - Inches(0.55), Inches(9), Inches(0.35),
        [[("ROOTS", HEAD, 11, c1, True),
          ("   \u00b7   S\u00e4ljkickoff 2026", BODY, 10, c2, False)]], space_after=0)


def new_slide(color=OFFWHITE):
    s = prs.slides.add_slide(BLANK)
    bg(s, color)
    return s


def section_head(slide, kick, title, y=Inches(0.72), tsize=32, tw=Inches(11.6)):
    kicker(slide, Inches(0.7), y, kick)
    txt(slide, Inches(0.66), y + Inches(0.42), tw, Inches(1.4),
        [[(title, HEAD, tsize, INK, True)]], space_after=0, line_spacing=1.03)


# ════════════════════════ SLIDE RENDERERS ════════════════════════
def r_cover(d):
    s = new_slide(OFFWHITE)
    img_cover(s, web_small("images/collection-2.jpg", 2000), Inches(7.7), 0, SW - Inches(7.7), SH)
    rect(s, Inches(7.7), 0, Inches(0.06), SH, FOREST)
    spine(s, FOREST)
    img_fit(s, brand("brand/roots-logo-black.png"),
            Inches(0.9), Inches(0.85), Inches(2.4), Inches(0.9), align="left")
    kicker(s, Inches(0.92), Inches(2.35), d["kicker"])
    txt(s, Inches(0.9), Inches(2.75), Inches(6.6), Inches(2.4),
        [[(d["title1"], HEAD, 40, INK, True)],
         [(d["title2"], HEAD, 40, FOREST, True)]],
        space_after=2, line_spacing=1.04)
    txt(s, Inches(0.92), Inches(4.85), Inches(6.2), Inches(1.4),
        [[(d["subtitle"], BODY, 15, SAND_DARK, False)]], line_spacing=1.35)


def r_agenda(d):
    s = new_slide(OFFWHITE)
    spine(s)
    section_head(s, d["kicker"], d["title"])
    items = d["items"]
    col_x = [Inches(0.72), Inches(6.9)]
    top = Inches(2.5)
    gap = Inches(0.92)
    per_col = 5
    for i, it in enumerate(items):
        col = 0 if i < per_col else 1
        row = i if i < per_col else i - per_col
        x = col_x[col]
        cy = top + row * gap
        txt(s, x, cy - Inches(0.02), Inches(0.7), Inches(0.6),
            [[(f"{i+1:02d}", HEAD, 20, OLIVE, True)]], space_after=0)
        txt(s, x + Inches(0.78), cy + Inches(0.03), Inches(5.2), Inches(0.7),
            [[(it, HEAD, 16, INK, True)]], space_after=0, line_spacing=1.05)
    footer(s); page_no(s)


def r_divider(d):
    s = new_slide(INK)
    spine(s, OLIVE)
    img_fit(s, brand("brand/roots-symbol-white.png"),
            SW - Inches(1.7), Inches(0.7), Inches(0.9), Inches(0.9), align="right")
    txt(s, Inches(0.9), Inches(2.4), Inches(6), Inches(1.6),
        [[(d["num"], HEAD, 84, OLIVE, True)]], space_after=0, line_spacing=1.0)
    txt(s, Inches(0.95), Inches(4.15), Inches(10.5), Inches(1.6),
        [[(d["title"], HEAD, 34, WHITE, True)]], space_after=0, line_spacing=1.05)
    rect(s, Inches(0.98), Inches(3.95), Inches(1.4), Inches(0.05), FOREST)
    footer(s, light=True); page_no(s, light=True)


def r_points(d):
    s = new_slide(OFFWHITE)
    spine(s)
    section_head(s, d["kicker"], d["title"])
    bullets = d["bullets"]
    n = len(bullets)
    top = Inches(2.7) if n <= 3 else Inches(2.5)
    gap = Inches(1.0) if n <= 3 else Inches(0.92)
    for i, b in enumerate(bullets):
        cy = top + i * gap
        dot(s, Inches(0.74), cy + Inches(0.09), Inches(0.2), FOREST)
        txt(s, Inches(1.2), cy, Inches(10.8), Inches(0.8),
            [[(b, HEAD, 20, INK, True)]], space_after=0, line_spacing=1.05)
    footer(s); page_no(s)


def r_hero(d):
    s = new_slide(FOREST_SOFT)
    spine(s, FOREST)
    kicker(s, Inches(0.9), Inches(2.2), d["kicker"], color=FOREST)
    txt(s, Inches(0.88), Inches(2.7), Inches(11.4), Inches(3.4),
        [[(d["line1"], HEAD, 40, INK, True)],
         [(d["line2"], HEAD, 40, FOREST, True)]],
        space_after=4, line_spacing=1.08)
    footer(s); page_no(s)


def _cards_geom(n):
    left = Inches(0.7)
    right = SW - Inches(0.7)
    gutter = Inches(0.35)
    total = right - left
    cw = int((total - gutter * (n - 1)) / n)
    xs = [left + i * (cw + gutter) for i in range(n)]
    return xs, cw


def r_cards(d):
    s = new_slide(OFFWHITE)
    spine(s)
    section_head(s, d["kicker"], d["title"])
    cards = d["cards"]
    n = len(cards)
    xs, cw = _cards_geom(n)
    cy = Inches(2.7)
    ch = Inches(3.4)
    for (t, b), x in zip(cards, xs):
        rrect(s, x, cy, cw, ch, WHITE, radius=0.05)
        rect(s, x + Inches(0.32), cy + Inches(0.42), Inches(0.5), Inches(0.06), FOREST)
        txt(s, x + Inches(0.32), cy + Inches(0.68), cw - Inches(0.64), Inches(0.9),
            [[(t, HEAD, 19, INK, True)]], space_after=0, line_spacing=1.05)
        txt(s, x + Inches(0.32), cy + Inches(1.55), cw - Inches(0.64), Inches(1.6),
            [[(b, BODY, 13, SAND_DARK, False)]], line_spacing=1.32)
    footer(s); page_no(s)


def r_products(d):
    s = new_slide(OFFWHITE)
    spine(s)
    section_head(s, d["kicker"], d["title"])
    imgs = ["images/schampoo.jpg", "images/conditioner.jpg", "images/body-wash.jpg"]
    cards = d["cards"]
    n = len(cards)
    xs, cw = _cards_geom(n)
    cy = Inches(2.55)
    ch = Inches(3.7)
    for (name, typ, benefit), x, im in zip(cards, xs, imgs):
        rrect(s, x, cy, cw, ch, WHITE, radius=0.05)
        img_cover(s, web_small(im, 1200), x, cy, cw, Inches(1.75))
        txt(s, x + Inches(0.3), cy + Inches(1.95), cw - Inches(0.6), Inches(0.5),
            [[(typ.upper(), BODY, 10.5, FOREST, True)]], space_after=0)
        txt(s, x + Inches(0.3), cy + Inches(2.3), cw - Inches(0.6), Inches(0.6),
            [[(name, HEAD, 20, INK, True)]], space_after=0, line_spacing=1.0)
        txt(s, x + Inches(0.3), cy + Inches(2.95), cw - Inches(0.6), Inches(0.7),
            [[(benefit, BODY, 12.5, SAND_DARK, False)]], line_spacing=1.3)
    footer(s); page_no(s)


def r_flow(d):
    s = new_slide(OFFWHITE)
    spine(s)
    section_head(s, d["kicker"], d["title"])
    steps = d["steps"]
    left = Inches(0.72)
    right = SW - Inches(0.72)
    per_row = 4
    gutter = Inches(0.3)
    total = right - left
    cw = int((total - gutter * (per_row - 1)) / per_row)
    ch = Inches(1.35)
    rows_y = [Inches(2.9), Inches(4.75)]
    for i, st in enumerate(steps):
        row = i // per_row
        col = i % per_row
        x = left + col * (cw + gutter)
        y = rows_y[row]
        rrect(s, x, y, cw, ch, FOREST_SOFT, radius=0.12)
        txt(s, x + Inches(0.28), y + Inches(0.2), cw - Inches(0.5), Inches(0.4),
            [[(f"{i+1:02d}", HEAD, 15, OLIVE, True)]], space_after=0)
        txt(s, x + Inches(0.28), y + Inches(0.62), cw - Inches(0.5), Inches(0.6),
            [[(st, HEAD, 15, INK, True)]], space_after=0, line_spacing=1.0)
        # arrow to next (same row)
        if col < per_row - 1 and i + 1 < len(steps):
            ax = x + cw + Inches(0.04)
            txt(s, ax, y + Inches(0.42), gutter - Inches(0.02), Inches(0.5),
                [[("\u203a", HEAD, 22, SAND_MED, True)]], align=PP_ALIGN.CENTER, space_after=0)
    footer(s); page_no(s)


def r_screens(d):
    s = new_slide(OFFWHITE)
    spine(s)
    section_head(s, d["kicker"], d["title"])
    screens = d["screens"]
    n = len(screens)
    xs, cw = _cards_geom(n)
    top = Inches(2.45)
    frame_h = Inches(3.4)
    for (img, cap, sub), x in zip(screens, xs):
        rrect(s, x, top, cw, frame_h, INK, radius=0.06)
        img_fit(s, web_small(img, 1200), x + Inches(0.12), top + Inches(0.12),
                cw - Inches(0.24), frame_h - Inches(0.24), valign="top")
        txt(s, x, top + frame_h + Inches(0.18), cw, Inches(0.5),
            [[(cap, HEAD, 16, INK, True)]], align=PP_ALIGN.CENTER, space_after=0)
        txt(s, x, top + frame_h + Inches(0.58), cw, Inches(0.5),
            [[(sub, BODY, 12, SAND_DARK, False)]], align=PP_ALIGN.CENTER, space_after=0)
    footer(s); page_no(s)


def r_close(d):
    s = new_slide(INK)
    spine(s, OLIVE)
    img_fit(s, brand("brand/roots-symbol-white.png"),
            Inches(0.9), Inches(0.85), Inches(0.95), Inches(0.95), align="left")
    txt(s, Inches(0.88), Inches(2.9), Inches(11.4), Inches(1.4),
        [[(d["title"], HEAD, 54, WHITE, True)]], space_after=0, line_spacing=1.0)
    txt(s, Inches(0.92), Inches(4.35), Inches(11), Inches(0.8),
        [[(d["subtitle"], BODY, 17, OLIVE, False)]], space_after=0)


RENDER = {
    "cover": r_cover, "agenda": r_agenda, "divider": r_divider, "points": r_points,
    "hero": r_hero, "cards": r_cards, "products": r_products, "flow": r_flow,
    "screens": r_screens, "close": r_close,
}


def title_for(d):
    k = d["kind"]
    if k == "cover":
        return "Titelsida"
    if k == "divider":
        return f'{d["num"]} \u00b7 {d["title"]}'
    if k == "hero":
        return f'{d["line1"]} {d["line2"]}'
    if k == "close":
        return d["title"]
    return d.get("title", "")


# ════════════════════════ BUILD PPTX ════════════════════════
for d in SLIDES:
    RENDER[d["kind"]](d)
prs.save(OUT_PPTX)

# Bädda in Alan Sans + Inter så filen renderas rätt på alla datorer.
import sys as _sys
_sys.path.insert(0, os.path.abspath(os.path.join(ROOT, "..", "roots-templates")))
try:
    from pptx_embed import embed_fonts
    embed_fonts(OUT_PPTX)
    print("PPTX:", OUT_PPTX, f"({len(SLIDES)} slides, typsnitt inbäddade)")
except Exception as _e:
    print("PPTX:", OUT_PPTX, f"({len(SLIDES)} slides) – font-embedding hoppades över:", _e)


# ════════════════════════ BUILD DOCX (talarmanus) ════════════════════════
DINK = DColor(0x1D, 0x1D, 0x1B)
DSAND = DColor(0x7F, 0x71, 0x5B)
DFOREST = DColor(0x6B, 0x79, 0x4F)
DOLIVE = DColor(0x9A, 0x98, 0x6B)


def set_font(run, name=BODY, size=11, color=DINK, bold=False, italic=False):
    run.font.name = name
    run.font.size = DPt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    for a in ("w:ascii", "w:hAnsi", "w:cs"):
        rfonts.set(qn(a), name)


def para(doc, runs, align=WD_ALIGN_PARAGRAPH.LEFT, space_after=6, space_before=0):
    p = doc.add_paragraph()
    p.alignment = align
    p.paragraph_format.space_after = DPt(space_after)
    p.paragraph_format.space_before = DPt(space_before)
    if isinstance(runs, tuple):
        runs = [runs]
    for seg in runs:
        text, kw = seg[0], (seg[1] if len(seg) > 1 else {})
        set_font(p.add_run(text), **kw)
    return p


def bullet(doc, text, color=DINK):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = DPt(4)
    p.paragraph_format.left_indent = DInches(0.28)
    p.paragraph_format.first_line_indent = DInches(-0.2)
    set_font(p.add_run("\u2013  "), color=DFOREST, bold=True)
    set_font(p.add_run(text), color=color, size=11.5)
    return p


doc = Document()
for sec in doc.sections:
    sec.top_margin = DInches(0.9); sec.bottom_margin = DInches(0.9)
    sec.left_margin = DInches(1.0); sec.right_margin = DInches(1.0)

# Title page
logo = brand("brand/roots-logo-black.png")
if logo:
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    p.add_run().add_picture(logo, width=DInches(1.9))
para(doc, [("S\u00e4ljkickoff 2026", {"name": HEAD, "size": 30, "color": DINK, "bold": True})],
     space_before=10, space_after=2)
para(doc, [("Talarmanus \u00b7 st\u00f6dord per slide", {"name": HEAD, "size": 15, "color": DFOREST, "bold": True})],
     space_after=10)
para(doc, [("Kort st\u00f6d till presentationen \u2013 s\u00e4g det med egna ord. Siffror inom [hakparentes] "
            "fylls i innan m\u00f6tet.", {"name": BODY, "size": 11.5, "color": DSAND, "italic": True})],
     space_after=4)
doc.add_page_break()

for i, d in enumerate(SLIDES, 1):
    para(doc, [(f"SLIDE {i:02d}", {"name": BODY, "size": 10, "color": DOLIVE, "bold": True})],
         space_before=(0 if i == 1 else 10), space_after=1)
    para(doc, [(title_for(d), {"name": HEAD, "size": 17, "color": DINK, "bold": True})], space_after=6)
    for note in d.get("notes", []):
        bullet(doc, note)

doc.save(OUT_DOCX)
print("DOCX:", OUT_DOCX)
