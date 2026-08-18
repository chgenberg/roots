"""Renderingsbibliotek för Roots-presentationer.

Nordisk minimalism enligt Roots brandbook: sandtoner, off-white bakgrund,
skogsgrön accent, Alan Sans som rubrikfont och Inter som brödtext. Generösa
marginaler, tunna linjer, inga skuggor och inga gradienter.

All geometri anges i tum som flyttal och konverteras till EMU i primitiverna,
så att layoutkoden går att läsa som en skiss.
"""

from __future__ import annotations

import os
from dataclasses import dataclass

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from slide_model import Slide

# ─────────────────────────────── Brand ────────────────────────────────
INK = RGBColor(0x1D, 0x1D, 0x1B)
SAND_DARK = RGBColor(0x7F, 0x71, 0x5B)
SAND_MED = RGBColor(0xB2, 0xA4, 0x91)
SAND_LIGHT = RGBColor(0xD5, 0xCA, 0xBF)
SAND_100 = RGBColor(0xF1, 0xEB, 0xE2)
OFFWHITE = RGBColor(0xFA, 0xF6, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOREST = RGBColor(0x6B, 0x79, 0x4F)
FOREST_SOFT = RGBColor(0xED, 0xF1, 0xE9)
OLIVE = RGBColor(0xC1, 0xBF, 0x99)

HEAD = "Alan Sans"
BODY = "Inter 18pt"

# ─────────────────────────────── Rutnät ───────────────────────────────
SW, SH = 13.333, 7.5
MARGIN = 0.85
CONTENT_W = SW - 2 * MARGIN
KICKER_Y = 0.60
TITLE_Y = 0.98
FOOTER_Y = 6.92

REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
PUBLIC = os.path.join(REPO, "apps", "web", "public")
# Bilder som bara hör till en presentation och inte till sajten.
ASSETS = os.path.join(os.path.dirname(__file__), "_assets")
CACHE = os.path.join(os.path.dirname(__file__), "_cache")


# ───────────────────────────── Bildhantering ──────────────────────────
def image_path(rel: str) -> str | None:
    """Skalar ned en bild ur apps/web/public eller _assets och returnerar den.

    Presentationerna mejlas runt, så originalen (flera megabyte styck) får inte
    bäddas in råa.
    """
    src = os.path.join(PUBLIC, rel)
    if not os.path.exists(src):
        src = os.path.join(ASSETS, rel)
    if not os.path.exists(src):
        return None
    os.makedirs(CACHE, exist_ok=True)
    dst = os.path.join(CACHE, rel.replace("/", "__"))
    if not os.path.exists(dst) or os.path.getmtime(dst) < os.path.getmtime(src):
        im = Image.open(src).convert("RGB")
        im.thumbnail((1700, 1700))
        im.save(dst, "JPEG", quality=82)
    return dst


def _content_bottom(im: Image.Image, skip_left: float) -> int:
    """Sista bildraden som innehåller något annat än bakgrund.

    Skärmbilder av portalen slutar ofta med en halv skärm tom yta. Sidomenyn
    går däremot ända ned, så vi tittar bara till höger om den — annars hittar
    vi alltid dess användarruta och kan aldrig beskära något.
    """
    w, h = im.size
    x0 = int(w * skip_left)
    strip = im.crop((x0, 0, w, h))
    bg = strip.getpixel((strip.width - 4, 4))
    px = strip.load()
    step = max(1, strip.width // 90)
    for y in range(h - 1, 0, -1):
        for x in range(0, strip.width, step):
            r, g, b = px[x, y]
            if abs(r - bg[0]) + abs(g - bg[1]) + abs(b - bg[2]) > 24:
                return y
    return h - 1


def screenshot_path(rel: str, skip_left: float = 0.0) -> str | None:
    """Beskär bort tom yta under innehållet och cachar skärmbilden.

    Sparas som PNG, inte JPEG: en skärmbild är gränssnittstext mot enfärgad
    botten, och där syns JPEG-artefakterna direkt runt bokstäverna.
    """
    src = os.path.join(ASSETS, "skarm", rel)
    if not os.path.exists(src):
        return None
    os.makedirs(CACHE, exist_ok=True)
    dst = os.path.join(CACHE, f"skarm__{rel}")
    if os.path.exists(dst) and os.path.getmtime(dst) >= os.path.getmtime(src):
        return dst

    im = Image.open(src).convert("RGB")
    bottom = _content_bottom(im, skip_left)
    # Lite luft under sista elementet, och beskär bara när det är värt det.
    pad = int(im.height * 0.03)
    cut = min(im.height, bottom + pad)
    if cut < im.height * 0.92:
        im = im.crop((0, 0, im.width, cut))
    im.thumbnail((1600, 1600))
    im.save(dst, "PNG", optimize=True)
    return dst


# ─────────────────────────────── Primitiver ───────────────────────────
def background(slide, color=OFFWHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _shape(slide, kind, x, y, w, h, fill, line=None, line_pt=1.0):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_pt)
    sp.shadow.inherit = False
    return sp


def rect(slide, x, y, w, h, fill, line=None, line_pt=1.0):
    return _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line, line_pt)


def rrect(slide, x, y, w, h, fill, line=None, radius=0.045, line_pt=1.0):
    sp = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line, line_pt)
    try:
        sp.adjustments[0] = radius
    except (IndexError, AttributeError):
        pass
    return sp


def oval(slide, x, y, w, h, fill, line=None, line_pt=1.0):
    return _shape(slide, MSO_SHAPE.OVAL, x, y, w, h, fill, line, line_pt)


def hairline(slide, x, y, w, color=SAND_LIGHT, thickness=0.011):
    return rect(slide, x, y, w, thickness, color)


def text(slide, x, y, w, h, paragraphs, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space_after=6, line_spacing=1.18):
    """Skriver en textbox där varje paragraf är en lista av segment.

    Ett segment är (text, font, storlek, färg) med valfri femte post för fet
    stil, så att en rad kan blanda vikter utan extra boxar.
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0

    if paragraphs and isinstance(paragraphs[0], tuple):
        paragraphs = [paragraphs]

    for i, segments in enumerate(paragraphs):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = align
        para.space_before = Pt(0)
        para.space_after = Pt(space_after)
        para.line_spacing = line_spacing
        for segment in segments:
            content, font, size, color = segment[:4]
            bold = segment[4] if len(segment) > 4 else False
            run = para.add_run()
            run.text = content
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
    return box


def image_cover(slide, path, x, y, w, h, rounded=False, focus_y=0.5):
    """Fyller rutan helt med bilden och beskär överskottet, som CSS object-fit.

    focus_y väljer vilken del av bilden som behålls när höjden beskärs: 0 är
    överkanten, 1 underkanten. Produktbilderna har flaskan i övre halvan, så de
    behöver ett lägre värde än mitten.
    """
    if not path or not os.path.exists(path):
        return None
    iw, ih = Image.open(path).size
    box_ratio, img_ratio = w / h, iw / ih
    if img_ratio > box_ratio:
        draw_h, draw_w = h, h * img_ratio
    else:
        draw_w, draw_h = w, w / img_ratio
    pic = slide.shapes.add_picture(path, Inches(x - (draw_w - w) / 2),
                                  Inches(y - (draw_h - h) / 2),
                                  Inches(draw_w), Inches(draw_h))
    pic.crop_left = pic.crop_right = max(0.0, (draw_w - w) / (2 * draw_w))
    overflow = max(0.0, (draw_h - h) / draw_h)
    pic.crop_top = overflow * focus_y
    pic.crop_bottom = overflow * (1 - focus_y)
    left, top = Inches(x), Inches(y)
    pic.left, pic.top, pic.width, pic.height = left, top, Inches(w), Inches(h)
    if rounded:
        pic.auto_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        try:
            pic.adjustments[0] = 0.03
        except (IndexError, AttributeError):
            pass
    return pic


def image_contain(slide, path, x, y, w, h):
    """Lägger bilden hel inuti rutan och centrerar den, utan beskärning.

    Används för bilder som inte tål beskärning — ett fotograferat föremål mot
    vit botten förlorar sin poäng om kanterna skärs bort.
    """
    if not path or not os.path.exists(path):
        return None
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    draw_w, draw_h = iw * scale, ih * scale
    return slide.shapes.add_picture(
        path, Inches(x + (w - draw_w) / 2), Inches(y + (h - draw_h) / 2),
        Inches(draw_w), Inches(draw_h))


# ───────────────────── Uppskattning av radbrytning ────────────────────
# Bredden per tecken som andel av fontstorleken. Uppmätt i renderad utskrift
# till cirka 0,46 för båda fonterna; värdena här ligger medvetet något högre så
# att uppskattningen fäller till för liten text snarare än för stor. Används
# bara för att välja fontstorlek, aldrig för exakt positionering.
_CHAR_RATIO = {HEAD: 0.475, BODY: 0.480}


def _ratio(content: str, font: str) -> float:
    """Teckenbredd justerad för versalandel.

    "Det här är ROOTS" tar betydligt mer plats än lika många gemener, och utan
    justeringen godkänns en rubrikstorlek som sedan bryts på två rader.
    """
    letters = [c for c in content if c.isalpha()]
    upper = sum(1 for c in letters if c.isupper()) / len(letters) if letters else 0
    return _CHAR_RATIO.get(font, 0.5) * (1 + 0.28 * upper)


def chars_per_line(content: str, width_in: float, size_pt: float, font: str) -> int:
    return max(6, int(width_in * 72 / (size_pt * _ratio(content, font))))


def line_count(content: str, width_in: float, size_pt: float, font: str = BODY) -> int:
    if not content:
        return 0
    # Explicit \n i rubriker måste räknas — annars krockar underrubrik med
    # titel när presentatören skrivit in manuella radbrytningar.
    total = 0
    for para in content.split("\n"):
        if not para:
            total += 1
            continue
        per_line = chars_per_line(para, width_in, size_pt, font)
        lines, used = 1, 0
        for word in para.split():
            needed = len(word) + (1 if used else 0)
            if used + needed > per_line:
                lines += 1
                used = len(word)
            else:
                used += needed
        total += lines
    return total


def fit_size(content: str, width_in: float, base_pt: float, min_pt: float,
             max_lines: int, font: str = BODY) -> float:
    """Största storlek mellan min_pt och base_pt som håller sig inom max_lines.

    Utöver radantalet måste det längsta ordet rymmas på en rad. Utan det
    villkoret godkänns storlekar där ett långt sammansatt ord som
    "föreningsförsäljning" bryts mitt i sig självt.
    """
    longest = max((len(w) for w in content.split()), default=0)
    size = base_pt
    while size > min_pt:
        per_line = chars_per_line(content, width_in, size, font)
        if line_count(content, width_in, size, font) <= max_lines and longest <= per_line:
            break
        size -= 0.5
    return size


# ──────────────────────────── Slidekomponenter ────────────────────────
def kicker(slide, content, x=MARGIN, y=KICKER_Y, color=SAND_DARK, w=CONTENT_W,
           align=PP_ALIGN.LEFT):
    if not content:
        return
    # Spärrad versal-kicker: mellanslag mellan tecknen ger den lugna,
    # typografiska känslan utan att kräva egna teckenavstånd i OOXML.
    spaced = " ".join(content.upper())
    text(slide, x, y, w, 0.3, [[(spaced, BODY, 9.5, color, True)]], align=align,
         space_after=0)


def title_block(slide, slide_spec, color=INK, sub_color=SAND_DARK,
                base_pt=37.0, width=None, gap=0.0) -> float:
    """Skriver kicker, rubrik, underrubrik och linje. Returnerar y för brödtext.

    gap skjuter ned linjen under rubriken. Extra luft efter rubrik så att
    underhäng (g, j, y) inte krockar med accentlinjen.
    """
    width = width or CONTENT_W
    kicker(slide, slide_spec.kicker)

    size = fit_size(slide_spec.title, width, base_pt, 21, 2, HEAD)
    lines = line_count(slide_spec.title, width, size, HEAD)
    # Generös höjd + luft: Alan Sans-underhäng (g/j/y) går långt under baseline.
    title_h = lines * size * 1.22 / 72
    text(slide, MARGIN, TITLE_Y, width, title_h + 0.14,
         [[(slide_spec.title, HEAD, size, color, False)]], space_after=0,
         line_spacing=1.08)
    y = TITLE_Y + title_h + 0.52

    if slide_spec.subtitle:
        sub_pt = fit_size(slide_spec.subtitle, width, 14.5, 11, 2)
        sub_h = line_count(slide_spec.subtitle, width, sub_pt) * sub_pt * 1.3 / 72
        text(slide, MARGIN, y, width, sub_h + 0.06,
             [[(slide_spec.subtitle, BODY, sub_pt, sub_color, False)]], space_after=0)
        y += sub_h + 0.28

    hairline(slide, MARGIN, y + gap, 1.5, FOREST, 0.022)
    return y + gap + 0.40


def footer(slide, deck_name, page, light=False):
    color = SAND_MED if not light else RGBColor(0xB9, 0xC0, 0xAC)
    text(slide, MARGIN, FOOTER_Y, 6.0, 0.3,
         [[(deck_name, BODY, 8.5, color, False)]], space_after=0)
    text(slide, SW - MARGIN - 1.2, FOOTER_Y, 1.2, 0.3,
         [[(f"{page:02d}", BODY, 8.5, color, False)]], align=PP_ALIGN.RIGHT,
         space_after=0)


def caption(slide, spec, y, w=CONTENT_W):
    """Master Sources "Nederst:"-rad — en sammanfattande rad under innehållet."""
    if not getattr(spec, "caption", ""):
        return
    size = fit_size(spec.caption, w, 14, 10.5, 2)
    text(slide, MARGIN, y, w, 0.62,
         [[(spec.caption, BODY, size, FOREST, False)]], space_after=0,
         line_spacing=1.3)


def bullets_extent(items, w, base_pt=14.5, gap_extra=0.20, max_h=None):
    """Storleken och höjden som bullets() skulle välja för samma indata.

    Finns för att en layout ska kunna centrera punktlistan innan den ritas —
    höjden beror på vilken storlek listan landar på, och den vet man annars
    först efteråt.
    """
    if not items:
        return base_pt, 0.0

    def height(size):
        return sum(line_count(i, w - 0.34, size) * size * 1.32 / 72 + gap_extra
                   for i in items)

    size = base_pt
    if max_h:
        while size > 10.5 and height(size) > max_h:
            size -= 0.5
    return size, height(size) - gap_extra


def bullets(slide, items, x, y, w, base_pt=14.5, gap_extra=0.20, dot_color=FOREST,
            color=INK, max_h=None):
    """Punktlista med små runda markörer. Krymper texten om utrymmet är knappt."""
    if not items:
        return y
    size, _ = bullets_extent(items, w, base_pt, gap_extra, max_h)

    dot = 0.075
    cursor = y
    for item in items:
        rows = line_count(item, w - 0.34, size)
        line_h = size * 1.32 / 72
        height = rows * line_h
        # Markören centreras på första radens mitt, inte på radens överkant.
        oval(slide, x + 0.02, cursor + line_h / 2 - dot / 2, dot, dot, dot_color)
        text(slide, x + 0.34, cursor, w - 0.34, height + 0.05,
             [[(item, BODY, size, color, False)]], space_after=0, line_spacing=1.32)
        cursor += height + gap_extra
    return cursor


# ─────────────────────────────── Layouter ─────────────────────────────
@dataclass
class Ctx:
    """Det renderarna behöver veta utöver själva sliden."""
    deck_name: str
    page: int


def r_cover(slide, spec: Slide, ctx: Ctx):
    has_image = bool(spec.images) and image_path(spec.images[0])
    # Luft mellan textpanel och foto — annars äter motivet i bilden
    # sista orden i rubriken visuellt.
    panel_w = 6.85 if has_image else 8.40
    gutter = 0.28 if has_image else 0.0

    background(slide, OFFWHITE)
    if has_image:
        # Lite mer av nedre halvan så en förgrundsprodukt inte äter hela panelens höjd.
        image_cover(slide, image_path(spec.images[0]), panel_w + gutter, 0,
                    SW - panel_w - gutter, SH, focus_y=0.58)
    else:
        # Utan foto håller ett sandfält med märket balansen i uppslaget.
        rect(slide, panel_w, 0, SW - panel_w, SH, SAND_100)
        symbol = os.path.join(PUBLIC, "brand", "roots-symbol-black.png")
        if os.path.exists(symbol):
            pic = slide.shapes.add_picture(symbol, Inches(0), Inches(0),
                                          height=Inches(1.45))
            pic.left = Inches(panel_w + (SW - panel_w) / 2) - pic.width // 2
            pic.top = Inches(SH / 2) - pic.height // 2

    logo = os.path.join(PUBLIC, "brand", "roots-logo-black.png")
    if os.path.exists(logo):
        slide.shapes.add_picture(logo, Inches(MARGIN), Inches(0.82), height=Inches(0.42))

    inner_w = panel_w - MARGIN - 0.55
    kicker(slide, spec.kicker, y=2.10, color=SAND_DARK, w=inner_w)

    size = fit_size(spec.title, inner_w, 44, 26, 5, HEAD)
    lines = line_count(spec.title, inner_w, size, HEAD)
    title_h = lines * size * 1.12 / 72
    text(slide, MARGIN, 2.50, inner_w, title_h + 0.12,
         [[(spec.title, HEAD, size, INK, False)]], space_after=0, line_spacing=1.08)

    y = 2.50 + title_h + 0.38
    hairline(slide, MARGIN, y, 1.5, FOREST, 0.022)
    if spec.subtitle:
        sub_pt = fit_size(spec.subtitle, inner_w, 14.5, 11.5, 4)
        text(slide, MARGIN, y + 0.32, inner_w, 1.15,
             [[(spec.subtitle, BODY, sub_pt, SAND_DARK, False)]], space_after=0,
             line_spacing=1.32)

    text(slide, MARGIN, FOOTER_Y, 6.0, 0.3,
         [[(ctx.deck_name, BODY, 9, SAND_MED, False)]], space_after=0)


def r_hero(slide, spec: Slide, ctx: Ctx):
    background(slide, INK)
    rect(slide, 0, 0, 0.10, SH, FOREST)

    kicker(slide, spec.kicker, x=1.30, y=1.55, color=OLIVE, w=CONTENT_W)

    width = SW - 1.30 - 1.30
    size = fit_size(spec.title, width, 44, 21, 4, HEAD)
    lines = line_count(spec.title, width, size, HEAD)
    title_h = lines * size * 1.16 / 72
    text(slide, 1.30, 2.05, width, title_h + 0.1,
         [[(spec.title, HEAD, size, WHITE, False)]], space_after=0, line_spacing=1.13)

    y = 2.05 + title_h + 0.34
    hairline(slide, 1.30, y, 1.5, OLIVE, 0.022)
    if spec.subtitle:
        sub_pt = fit_size(spec.subtitle, width - 1.0, 16, 12, 3)
        text(slide, 1.30, y + 0.40, width - 1.0, 1.2,
             [[(spec.subtitle, BODY, sub_pt, SAND_LIGHT, False)]], space_after=0)

    footer(slide, ctx.deck_name, ctx.page, light=True)


def r_points(slide, spec: Slide, ctx: Ctx):
    background(slide)
    img = image_path(spec.images[0]) if spec.images else None
    body_w = 6.55 if img else CONTENT_W

    y = title_block(slide, spec, width=body_w if img else CONTENT_W)
    if img:
        # Kortare bild när caption finns — annars krockar nedersta raden
        # med fotots nedre hörn.
        img_h = 4.85 if spec.caption else 5.35
        image_cover(slide, img, SW - MARGIN - 4.20, 1.10, 4.20, img_h,
                    rounded=True)

    cap_room = 0.78 if spec.caption else 0
    bullets(slide, spec.items, MARGIN, y, body_w,
            max_h=FOOTER_Y - 0.35 - y - cap_room)
    if spec.caption:
        caption(slide, spec, FOOTER_Y - 0.82, w=body_w)
    footer(slide, ctx.deck_name, ctx.page)


def r_cards(slide, spec: Slide, ctx: Ctx):
    background(slide)
    y = title_block(slide, spec)

    rows = spec.cells(2)
    n = max(1, len(rows))
    gap = 0.30
    # Fem kort på en rad blir så smala att långa ord bryts mitt itu, och fyra
    # blir för trånga när brödtexten är lång. Då är flera rader läsbarare.
    if n >= 5:
        # Med lång brödtext behöver korten bredd mer än de behöver ligga på
        # samma rad — annars bryts både rubriker och text mitt itu.
        per_row = 2 if max(len(b) for _, b in rows) > 70 else 3
    elif n == 4 and max(len(b) for _, b in rows) > 70:
        per_row = 2
    else:
        per_row = n
    row_count = (n + per_row - 1) // per_row
    card_w = (CONTENT_W - gap * (per_row - 1)) / per_row
    bottom = FOOTER_Y - 0.35 - (0.62 if spec.caption else 0)
    card_h = min(3.35, (bottom - y - gap * (row_count - 1)) / row_count)

    # Ett lågt och brett kort får rubrik och text sida vid sida — staplat
    # vertikalt skulle texten tryckas ut genom kortets underkant.
    wide = card_h < 2.4

    if wide:
        # Rubrikspalten måste rymma det längsta ordet — "Organisationsutveckling"
        # och "Genomförandekraft" bryts mitt itu om spalten är för smal.
        # fit_size krymper bara till min_pt och ger upp där, så bredden måste
        # växa i stället.
        head_w = card_w * 0.34
        longest = max((w for h, _ in rows for w in h.split()), key=len)
        while (head_w < card_w * 0.46
               and chars_per_line(longest, head_w, 12, HEAD) < len(longest) + 2):
            head_w += 0.04
        body_w = card_w - head_w - 0.98
        # fit_size godkänner storlekar där det längsta ordet nätt och jämnt
        # ryms och tar inte hänsyn till kortets höjd. Båda marginalerna behövs:
        # teckenbredden är en uppskattning, och en tvåradig rubrik som växer ur
        # sin box lägger sig över numret ovanför.
        head_room = max(0.40, card_h - 0.90)
        head_pt = min(fit_size(h, head_w, 16, 11, 2, HEAD) for h, _ in rows)
        while head_pt > 10 and (
            chars_per_line(longest, head_w, head_pt, HEAD) < len(longest) + 1
            or max(line_count(h, head_w, head_pt, HEAD) for h, _ in rows)
            * head_pt * 1.16 / 72 > head_room
        ):
            head_pt -= 0.5
        # Texten är vertikalt centrerad i kortet, så för många rader spiller
        # ut både över och under kanten. Krymp till kortets faktiska höjd.
        body_room = card_h - 0.50
        body_pt = 12.5
        while body_pt > 8.5 and any(
            line_count(b, body_w, body_pt) * body_pt * 1.34 / 72 > body_room
            for _, b in rows
        ):
            body_pt -= 0.5
    else:
        inner = card_w - 0.84
        # Gemensam rubrikhöjd så att kortens brödtext börjar på samma rad.
        head_pt = min(fit_size(h, inner, 17, 12.5, 2, HEAD) for h, _ in rows)
        head_lines = max(line_count(h, inner, head_pt, HEAD) for h, _ in rows)
        head_h = head_lines * head_pt * 1.2 / 72
        body_offset = 0.92 + head_h + 0.18
        body_room = card_h - 0.40 - body_offset

    for i, (heading, body) in enumerate(rows):
        x = MARGIN + (i % per_row) * (card_w + gap)
        top = y + (i // per_row) * (card_h + gap)
        rrect(slide, x, top, card_w, card_h, WHITE, SAND_LIGHT)
        rect(slide, x, top, card_w, 0.035, FOREST)

        if wide:
            text(slide, x + 0.42, top + 0.30, 0.5, 0.26,
                 [[(f"{i + 1:02d}", BODY, 9.5, FOREST, True)]], space_after=0)
            # Rubriken centreras i den yta som blir kvar under numret. Startade
            # den vid numrets överkant växte en tvåradig rubrik uppåt och lade
            # sig över siffran.
            text(slide, x + 0.42, top + 0.58, head_w, card_h - 0.86,
                 [[(heading, HEAD, head_pt, INK, False)]], space_after=0,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.16)
            rect(slide, x + 0.42 + head_w + 0.28, top + 0.42, 0.011,
                 card_h - 0.84, SAND_LIGHT)
            text(slide, x + 0.42 + head_w + 0.42, top + 0.30, body_w,
                 card_h - 0.5,
                 [[(body, BODY, body_pt, SAND_DARK, False)]], space_after=0,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.34)
            continue

        text(slide, x + 0.42, top + 0.46, 0.5, 0.3,
             [[(f"{i + 1:02d}", BODY, 10, FOREST, True)]], space_after=0)
        text(slide, x + 0.42, top + 0.86, inner, head_h + 0.08,
             [[(heading, HEAD, head_pt, INK, False)]], space_after=0, line_spacing=1.16)

        body_pt = 13.0
        while body_pt > 9.5 and (
            line_count(body, inner, body_pt) * body_pt * 1.34 / 72 > body_room
        ):
            body_pt -= 0.5
        text(slide, x + 0.42, top + body_offset, inner, body_room,
             [[(body, BODY, body_pt, SAND_DARK, False)]], space_after=0,
             line_spacing=1.34)

    footer(slide, ctx.deck_name, ctx.page)
    caption(slide, spec, y + row_count * card_h + (row_count - 1) * gap + 0.22)


def r_columns(slide, spec: Slide, ctx: Ctx):
    background(slide)
    y = title_block(slide, spec)

    gap = 0.42
    panel_w = (CONTENT_W - gap) / 2
    panel_h = min(3.75, FOOTER_Y - 0.35 - y - (0.62 if spec.caption else 0))

    for i, (column, fill, accent) in enumerate((
        (spec.left, SAND_100, SAND_MED),
        (spec.right, FOREST_SOFT, FOREST),
    )):
        if column is None:
            continue
        x = MARGIN + i * (panel_w + gap)
        rrect(slide, x, y, panel_w, panel_h, fill)
        rect(slide, x, y, 0.035, panel_h, accent)

        head_pt = fit_size(column.heading, panel_w - 1.0, 19, 14, 2, HEAD)
        head_h = line_count(column.heading, panel_w - 1.0, head_pt, HEAD) * head_pt * 1.2 / 72
        text(slide, x + 0.50, y + 0.44, panel_w - 1.0, head_h + 0.1,
             [[(column.heading, HEAD, head_pt, INK, False)]], space_after=0,
             line_spacing=1.16)

        top = y + 0.50 + head_h + 0.26
        bullets(slide, column.items, x + 0.50, top, panel_w - 1.0, base_pt=13.5,
                gap_extra=0.16, dot_color=accent, max_h=y + panel_h - 0.40 - top)

    footer(slide, ctx.deck_name, ctx.page)
    caption(slide, spec, y + panel_h + 0.22)


def r_timeline(slide, spec: Slide, ctx: Ctx):
    background(slide)
    y = title_block(slide, spec)

    rows = spec.cells(2)
    n = max(1, len(rows))
    col_w = CONTENT_W / n
    rail_y = y + 0.55

    hairline(slide, MARGIN + col_w / 2, rail_y, CONTENT_W - col_w, SAND_LIGHT, 0.014)

    for i, (label, body) in enumerate(rows):
        cx = MARGIN + col_w * i + col_w / 2
        oval(slide, cx - 0.17, rail_y - 0.16, 0.34, 0.34, OFFWHITE, FOREST, 1.25)
        text(slide, cx - 0.17, rail_y - 0.095, 0.34, 0.24,
             [[(str(i + 1), BODY, 10, FOREST, True)]], align=PP_ALIGN.CENTER,
             space_after=0)

        inner = col_w - 0.42
        label_pt = fit_size(label, inner, 16, 11.5, 2, HEAD)
        label_h = line_count(label, inner, label_pt, HEAD) * label_pt * 1.2 / 72
        text(slide, cx - inner / 2, rail_y + 0.46, inner, label_h + 0.08,
             [[(label, HEAD, label_pt, INK, False)]], align=PP_ALIGN.CENTER,
             space_after=0, line_spacing=1.16)

        body_pt = fit_size(body, inner, 12.5, 10, 6)
        text(slide, cx - inner / 2, rail_y + 0.52 + label_h + 0.14, inner, 1.9,
             [[(body, BODY, body_pt, SAND_DARK, False)]], align=PP_ALIGN.CENTER,
             space_after=0, line_spacing=1.34)

    footer(slide, ctx.deck_name, ctx.page)
    caption(slide, spec, rail_y + 2.7)


def r_quarters(slide, spec: Slide, ctx: Ctx):
    background(slide)
    y = title_block(slide, spec)

    rows = spec.cells(3)
    n = max(1, len(rows))
    gap = 0.26
    card_w = (CONTENT_W - gap * (n - 1)) / n
    # Korten är smala, så fokustexten behöver höjd: Q4:s fokusmening tar fyra
    # rader. Med ett lägre tak lämnades en tom remsa under korten samtidigt som
    # texten trängdes ihop inuti dem.
    card_h = min(4.40, FOOTER_Y - 0.35 - y - (0.72 if spec.caption else 0))

    # Kvartalsbeteckningen bryts ut ur etiketten så att "Q1" kan bli en chip.
    # Master Source skriver dem som "Q1 – Commercial Proof of Concept".
    split = []
    for label in (r[0] for r in rows):
        head, _, rest = label.partition(" – ")
        split.append((head, rest) if rest else tuple((label.split(" ", 1) + [""])[:2]))
    inner = card_w - 0.68

    # Rubrikerna är olika långa ("Operational Readiness" mot "Commercial
    # Execution & Delivery"). Gemensam storlek och höjd gör att fokustexten och
    # leveransfältet linjerar mellan korten. På ett lågt kort får rubriken bara
    # ta två rader, annars äter den upp fokustexten.
    # Fyra kort blir så smala att en trerading rubrik ("Commercial Execution &
    # Delivery") äter upp den höjd fokustexten behöver längre ned i kortet.
    name_max = 3 if card_h >= 3.5 and n <= 3 else 2
    name_pt = min(fit_size(name, inner, 16, 11.5, name_max, HEAD) for _, name in split)
    name_lines = max(line_count(name, inner, name_pt, HEAD) for _, name in split)
    name_h = name_lines * name_pt * 1.2 / 72

    # Leveransfältet ligger i botten och får den höjd dess text kräver, men
    # aldrig så mycket att fokustexten trycks bort.
    del_pt = min(fit_size(d, inner, 12, 9.5, 4) for _, _, d in rows)
    del_lines = min(4, max(line_count(d, inner, del_pt) for _, _, d in rows))
    del_h = del_lines * del_pt * 1.30 / 72
    foot_h = 0.20 + del_h + 0.30

    focus_top = y + 0.98 + name_h + 0.20
    focus_room = max(0.42, (y + card_h - foot_h - 0.14) - focus_top)

    for i, (label, focus, delivery) in enumerate(rows):
        x = MARGIN + i * (card_w + gap)
        rrect(slide, x, y, card_w, card_h, WHITE, SAND_LIGHT)
        rect(slide, x, y, card_w, 0.035, FOREST)

        quarter, name = split[i]
        rrect(slide, x + 0.34, y + 0.40, 0.62, 0.34, FOREST_SOFT, radius=0.30)
        text(slide, x + 0.34, y + 0.465, 0.62, 0.26,
             [[(quarter, BODY, 11, FOREST, True)]], align=PP_ALIGN.CENTER,
             space_after=0)

        text(slide, x + 0.34, y + 0.94, inner, name_h + 0.08,
             [[(name, HEAD, name_pt, INK, False)]], space_after=0, line_spacing=1.16)

        # Radberäkningen är en uppskattning och hamnar en rad fel på de längsta
        # fokusmeningarna. Marginalen håller texten ovanför leveransfältet.
        focus_pt = 12.5
        while focus_pt > 9.0 and (
            line_count(focus, inner * 0.94, focus_pt) * focus_pt * 1.36 / 72
            > focus_room - 0.06
        ):
            focus_pt -= 0.5
        text(slide, x + 0.34, focus_top, inner, focus_room,
             [[(focus, BODY, focus_pt, SAND_DARK, False)]], space_after=0,
             line_spacing=1.36)

        hairline(slide, x + 0.34, y + card_h - foot_h, inner, SAND_LIGHT)
        text(slide, x + 0.34, y + card_h - foot_h + 0.14, inner, 0.24,
             [[("LEVERANS", BODY, 8, SAND_MED, True)]], space_after=0)
        text(slide, x + 0.34, y + card_h - foot_h + 0.40, inner, del_h + 0.1,
             [[(delivery, BODY, del_pt, INK, False)]], space_after=0,
             line_spacing=1.30)

    footer(slide, ctx.deck_name, ctx.page)
    caption(slide, spec, y + card_h + 0.24)


def r_years(slide, spec: Slide, ctx: Ctx):
    background(slide)
    y = title_block(slide, spec)

    rows = spec.cells(3)
    n = max(1, len(rows))
    available = FOOTER_Y - 0.32 - y
    row_h = min(0.86, available / n)

    for i, (label, body, tag) in enumerate(rows):
        top = y + i * row_h
        if i % 2 == 0:
            rect(slide, MARGIN, top, CONTENT_W, row_h - 0.06, SAND_100)

        parts = label.split(" ", 2)
        year = " ".join(parts[:2]) if len(parts) > 1 else label
        name = parts[2] if len(parts) > 2 else ""

        # Alla fyra spalter centreras kring bandets mittlinje så att årtal,
        # namn, beskrivning och nyckelord ligger på samma optiska rad.
        mid = top + (row_h - 0.06) / 2

        text(slide, MARGIN + 0.30, mid - 0.15, 0.95, 0.3,
             [[(year.upper(), BODY, 10, FOREST, True)]], space_after=0,
             anchor=MSO_ANCHOR.MIDDLE)

        name_pt = fit_size(name, 2.95, 15.5, 11, 2, HEAD)
        text(slide, MARGIN + 1.32, mid - 0.25, 2.95, 0.5,
             [[(name, HEAD, name_pt, INK, False)]], space_after=0,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.14)

        body_pt = fit_size(body, 4.25, 12.5, 9.5, 2)
        text(slide, MARGIN + 4.42, mid - 0.275, 4.25, 0.55,
             [[(body, BODY, body_pt, INK, False)]], space_after=0,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.28)

        tag_pt = fit_size(tag, 2.55, 11, 8.5, 3)
        text(slide, SW - MARGIN - 2.60, mid - 0.30, 2.55, 0.6,
             [[(tag, BODY, tag_pt, SAND_DARK, False)]], align=PP_ALIGN.RIGHT,
             space_after=0, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.26)

    footer(slide, ctx.deck_name, ctx.page)


def r_flywheel(slide, spec: Slide, ctx: Ctx):
    """Ring av noder runt en mittcirkel.

    Bara nodnamnen sätts på sliden — beskrivningarna hör i talarmanuset,
    annars blir ringen oläslig.
    """
    import math

    background(slide)
    kicker(slide, spec.kicker)
    size = fit_size(spec.title, CONTENT_W, 34, 22, 1, HEAD)
    text(slide, MARGIN, TITLE_Y, CONTENT_W, 0.6,
         [[(spec.title, HEAD, size, INK, False)]], space_after=0)
    if spec.subtitle:
        text(slide, MARGIN, TITLE_Y + 0.58, CONTENT_W, 0.35,
             [[(spec.subtitle, BODY, 13, SAND_DARK, False)]], space_after=0)

    rows = spec.cells(2)
    nodes = [r[0] for r in rows] or ["—"]
    # Fler noder behöver större ring, annars kolliderar etiketterna. Ringen
    # centreras i det som återstår mellan rubriken och sidfoten, och lämnar
    # plats för en avslutande rad när sliden har en.
    caption_h = 0.62 if spec.caption else 0
    top, bottom = 2.16, FOOTER_Y - 0.30 - caption_h
    radius = min(1.98, 1.12 + 0.085 * len(nodes), (bottom - top) / 2 - 0.24)
    cx, cy = SW / 2, (top + bottom) / 2

    # Ringen går genom noderna, så att cykeln läses som en cykel.
    oval(slide, cx - radius, cy - radius, 2 * radius, 2 * radius,
         None, SAND_LIGHT, 0.75)
    oval(slide, cx - 0.86, cy - 0.86, 1.72, 1.72, FOREST_SOFT)
    # Mittcirkelns text sätts per slide: MS-014 har till exempel ett
    # "Support Center" i navet, inte ROOTS-flywheelet.
    hub = spec.columns or ["ROOTS", "Flywheel"]
    hub_lines = [[(hub[0], HEAD, 15, FOREST, False)]]
    if len(hub) > 1 and hub[1]:
        hub_lines.append([(hub[1], BODY, 11, SAND_DARK, False)])
    text(slide, cx - 0.78, cy - 0.34, 1.56, 0.7, hub_lines,
         align=PP_ALIGN.CENTER, space_after=1)

    label_w = 2.62
    for i, name in enumerate(nodes):
        angle = -math.pi / 2 + i * 2 * math.pi / len(nodes)
        dx, dy = math.cos(angle), math.sin(angle)
        nx, ny = cx + dx * radius, cy + dy * radius

        oval(slide, nx - 0.13, ny - 0.13, 0.26, 0.26, OFFWHITE, FOREST, 1.25)
        text(slide, nx - 0.13, ny - 0.075, 0.26, 0.2,
             [[(str(i + 1), BODY, 8.5, FOREST, True)]], align=PP_ALIGN.CENTER,
             space_after=0)

        pt = fit_size(name, label_w - 0.2, 12, 9, 2, HEAD)
        rows_needed = line_count(name, label_w - 0.2, pt, HEAD)
        height = rows_needed * pt * 1.2 / 72 + 0.06

        if abs(dx) < 0.35:  # rakt över eller under ringen
            lx, align = cx - label_w / 2, PP_ALIGN.CENTER
            ly = ny - 0.30 - height if dy < 0 else ny + 0.30
        elif dx > 0:
            lx, align = nx + 0.30, PP_ALIGN.LEFT
            ly = ny - height / 2
        else:
            lx, align = nx - 0.30 - label_w, PP_ALIGN.RIGHT
            ly = ny - height / 2

        text(slide, lx, ly, label_w, height,
             [[(name, HEAD, pt, INK, False)]], align=align, space_after=0,
             line_spacing=1.18)

    footer(slide, ctx.deck_name, ctx.page)
    caption(slide, spec, min(cy + radius + 0.44, FOOTER_Y - 0.56))


def r_products(slide, spec: Slide, ctx: Ctx):
    background(slide)
    y = title_block(slide, spec)

    rows = spec.cells(2)
    n = max(1, len(rows))
    gap = 0.34
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = min(3.85, FOOTER_Y - 0.35 - y - (0.62 if spec.caption else 0))
    # Korta kort får ge bilden mindre plats, annars trycks brödtexten ut nedåt.
    img_h = card_h * (0.52 if card_h >= 3.3 else 0.44)

    # Bildfält kan innehålla en bild per produkt separerade med mittpunkt.
    images = spec.images if len(spec.images) == n else [None] * n

    # Gemensam rubrikstorlek så att brödtexten börjar på samma höjd i alla kort.
    inner = card_w - 0.72
    name_pt = min(fit_size(name, inner, 18, 13, 2, HEAD) for name, _ in rows)
    name_lines = max(line_count(name, inner, name_pt, HEAD) for name, _ in rows)
    name_h = name_lines * name_pt * 1.2 / 72

    body_top = img_h + 0.30 + name_h + 0.14
    body_room = max(0.4, card_h - body_top - 0.30)
    body_pt = 12.5
    while body_pt > 9.5 and any(
        line_count(body, inner, body_pt) * body_pt * 1.34 / 72 > body_room
        for _, body in rows
    ):
        body_pt -= 0.5

    for i, (name, body) in enumerate(rows):
        x = MARGIN + i * (card_w + gap)
        rrect(slide, x, y, card_w, card_h, WHITE, SAND_LIGHT)

        path = image_path(images[i]) if images[i] else None
        if path:
            # Landskapsbilder med hela flaskan: centrera beskärningen så locket
            # inte blir det enda som syns i det breda kortet.
            image_cover(slide, path, x + 0.02, y + 0.02, card_w - 0.04, img_h,
                        focus_y=0.42)
        else:
            rect(slide, x + 0.02, y + 0.02, card_w - 0.04, img_h, SAND_100)

        text(slide, x + 0.36, y + img_h + 0.30, inner, name_h + 0.08,
             [[(name, HEAD, name_pt, INK, False)]], space_after=0, line_spacing=1.16)

        text(slide, x + 0.36, y + body_top, inner, body_room,
             [[(body, BODY, body_pt, SAND_DARK, False)]], space_after=0,
             line_spacing=1.34)

    footer(slide, ctx.deck_name, ctx.page)
    caption(slide, spec, y + card_h + 0.24)


def r_table(slide, spec: Slide, ctx: Ctx):
    background(slide)
    y = title_block(slide, spec)

    rows = spec.cells(2)
    headers = spec.columns or ["", ""]
    left_w = CONTENT_W * 0.40
    right_w = CONTENT_W - left_w
    head_h = 0.36
    available = FOOTER_Y - 0.38 - y - head_h - (0.62 if spec.caption else 0)
    # Radhöjd får aldrig överskrida tillgängligt utrymme (summa/caption krockar).
    row_h = min(0.62, available / max(1, len(rows)))
    max_lines = 2 if row_h >= 0.52 else 1

    rect(slide, MARGIN, y, CONTENT_W, head_h, SAND_100)
    for i, (label, width) in enumerate(zip(headers[:2], (left_w, right_w))):
        x = MARGIN + (0 if i == 0 else left_w)
        text(slide, x + 0.28, y + 0.08, width - 0.4, 0.22,
             [[(label.upper(), BODY, 9, SAND_DARK, True)]], space_after=0)

    for i, (left, right) in enumerate(rows):
        top = y + head_h + i * row_h
        if i % 2 == 1:
            rect(slide, MARGIN, top, CONTENT_W, row_h, SAND_100)

        left_pt = fit_size(left, left_w - 0.55, 12.5, 9, max_lines)
        text(slide, MARGIN + 0.28, top + 0.03, left_w - 0.55, row_h - 0.06,
             [[(left, BODY, left_pt, INK, False)]], space_after=0,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)

        right_pt = fit_size(right, right_w - 0.55, 12, 9, max_lines)
        text(slide, MARGIN + left_w + 0.28, top + 0.03, right_w - 0.55,
             row_h - 0.06,
             [[(right, BODY, right_pt, SAND_DARK, False)]], space_after=0,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)

    footer(slide, ctx.deck_name, ctx.page)
    if spec.caption:
        caption(slide, spec, FOOTER_Y - 0.52)


def r_kpi(slide, spec: Slide, ctx: Ctx):
    """KPI-kort med fasta band — värde, linje, etikett — utan överlapp."""
    background(slide)
    y = title_block(slide, spec)

    rows = spec.cells(2)
    n = max(1, len(rows))
    per_row = n if n <= 4 else 3
    gap = 0.28
    tile_w = (CONTENT_W - gap * (per_row - 1)) / per_row
    row_count = (n + per_row - 1) // per_row
    available = FOOTER_Y - 0.40 - y - (0.70 if spec.caption else 0)
    tile_h = min(2.55, (available - gap * (row_count - 1)) / row_count)
    # Gemensam värdezon så linjen ligger lika i alla kort och aldrig skär text.
    value_band = min(1.05, max(0.78, tile_h * 0.38))

    for i, (value, label) in enumerate(rows):
        col, row = i % per_row, i // per_row
        x = MARGIN + col * (tile_w + gap)
        top = y + row * (tile_h + gap)

        rrect(slide, x, top, tile_w, tile_h, WHITE, SAND_LIGHT)
        inner = tile_w - 0.64
        vx = x + 0.32

        # Preferera en rad — två rader bara om värdet absolut inte får plats.
        value_pt = fit_size(value, inner, 26, 13, 1, HEAD)
        if line_count(value, inner, value_pt, HEAD) > 1:
            value_pt = fit_size(value, inner, 22, 12, 2, HEAD)
        text(slide, vx, top + 0.18, inner, value_band - 0.22,
             [[(value, HEAD, value_pt, FOREST, False)]], space_after=0,
             line_spacing=1.05, anchor=MSO_ANCHOR.MIDDLE)

        line_y = top + value_band
        hairline(slide, vx, line_y, 0.70, SAND_LIGHT, 0.016)

        label_top = line_y + 0.20
        label_h = top + tile_h - label_top - 0.22
        label_pt = fit_size(label, inner, 12, 9, 4)
        text(slide, vx, label_top, inner, max(0.35, label_h),
             [[(label, BODY, label_pt, SAND_DARK, False)]], space_after=0,
             line_spacing=1.28)

    footer(slide, ctx.deck_name, ctx.page)
    if spec.caption:
        caption(slide, spec,
                y + row_count * tile_h + (row_count - 1) * gap + 0.26)


def r_roles(slide, spec: Slide, ctx: Ctx):
    background(slide)
    y = title_block(slide, spec)

    rows = spec.cells(2)
    n = max(1, len(rows))
    per_row = 2 if n <= 4 else 3
    gap_x, gap_y = 0.34, 0.24
    tile_w = (CONTENT_W - gap_x * (per_row - 1)) / per_row
    row_count = (n + per_row - 1) // per_row
    available = FOOTER_Y - 0.35 - y
    tile_h = min(1.60, (available - gap_y * (row_count - 1)) / row_count)

    # Gemensam rubrikhöjd så att alla rutors brödtext börjar på samma rad.
    inner = tile_w - 0.5
    name_pt = min(fit_size(name, inner, 15.5, 11.5, 2, HEAD) for name, _ in rows)
    name_h = max(line_count(name, inner, name_pt, HEAD) for name, _ in rows)
    name_h *= name_pt * 1.2 / 72
    duty_room = tile_h - name_h - 0.16

    for i, (name, duty) in enumerate(rows):
        col, row = i % per_row, i // per_row
        x = MARGIN + col * (tile_w + gap_x)
        top = y + row * (tile_h + gap_y)

        rect(slide, x, top, 0.032, tile_h, FOREST)
        text(slide, x + 0.28, top + 0.02, inner, name_h + 0.08,
             [[(name, HEAD, name_pt, INK, False)]], space_after=0, line_spacing=1.16)

        # Ansvarslistorna är olika långa; texten krymps till rutans utrymme i
        # stället för att tillåtas växa ner i raden under.
        duty_pt = 12.5
        while duty_pt > 8.0 and (
            line_count(duty, inner, duty_pt) * duty_pt * 1.32 / 72 > duty_room
        ):
            duty_pt -= 0.5
        text(slide, x + 0.28, top + name_h + 0.14, inner, duty_room,
             [[(duty, BODY, duty_pt, SAND_DARK, False)]], space_after=0,
             line_spacing=1.32)

    footer(slide, ctx.deck_name, ctx.page)


def r_compare(slide, spec: Slide, ctx: Ctx):
    """Produktjämförelse: tre kategorirader × fyra kort (Roots + konkurrenter).

    Varje `items`-rad: kategori | varumärke | produkt | pris/100ml | förp | badge | feature
    Badge «Vårt pris» markerar Roots-kortet (grön bakgrund).
    """
    background(slide)
    y = title_block(slide, spec, base_pt=22, gap=0.0)

    rows_raw = [p for p in spec.items if p.strip()]
    parsed: list[tuple[str, str, str, str, str, str, str]] = []
    for item in rows_raw:
        parts = [x.strip() for x in item.split("|")]
        while len(parts) < 7:
            parts.append("")
        parsed.append(tuple(parts[:7]))  # type: ignore[arg-type]

    categories: list[tuple[str, list]] = []
    for row in parsed:
        cat = row[0]
        if not categories or categories[-1][0] != cat:
            categories.append((cat, [row]))
        else:
            categories[-1][1].append(row)

    # Ingen caption-zon — underrubriken bär förklaringen och korten får höjden.
    floor = FOOTER_Y - 0.22
    n_cat = max(1, len(categories))
    gap_y = 0.14
    avail = floor - y - gap_y * (n_cat - 1)
    block_h = avail / n_cat
    label_h = 0.17
    card_h = block_h - label_h - 0.02

    cols = 4
    gap_x = 0.12
    card_w = (CONTENT_W - gap_x * (cols - 1)) / cols

    for ci, (cat, products) in enumerate(categories):
        top = y + ci * (block_h + gap_y)
        pill_w = min(1.70, 0.18 * len(cat) + 0.70)
        rrect(slide, MARGIN, top, pill_w, 0.17, SAND_100, None, radius=0.35)
        text(slide, MARGIN, top + 0.005, pill_w, 0.16,
             [[(cat.upper(), BODY, 8, SAND_DARK, True)]],
             align=PP_ALIGN.CENTER, space_after=0)

        card_top = top + label_h
        for i, (_, brand, product, price, pack, badge, feature) in enumerate(
            products[:cols]
        ):
            x = MARGIN + i * (card_w + gap_x)
            ours = badge.lower().startswith("vårt") or brand.lower() == "roots"
            fill = FOREST_SOFT if ours else WHITE
            border = FOREST if ours else SAND_LIGHT
            rrect(slide, x, card_top, card_w, card_h, fill, border,
                  radius=0.035, line_pt=1.2 if ours else 0.85)

            pad_x = 0.10
            inner = card_w - 2 * pad_x
            # Relativa band (andel av kortets höjd) så allt ryms.
            def band(a: float, b: float) -> tuple[float, float]:
                return card_top + card_h * a, card_h * (b - a)

            by, bh = band(0.03, 0.15)
            brand_pt = fit_size(brand, inner, 10.5, 7.5, 1, HEAD)
            text(slide, x + pad_x, by, inner, bh,
                 [[(brand, HEAD, brand_pt, INK, False)]],
                 align=PP_ALIGN.CENTER, space_after=0, anchor=MSO_ANCHOR.MIDDLE)

            py, ph = band(0.15, 0.32)
            prod_pt = fit_size(product, inner, 8, 6.5, 2)
            text(slide, x + pad_x, py, inner, ph,
                 [[(product, BODY, prod_pt, SAND_DARK, False)]],
                 align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.05,
                 anchor=MSO_ANCHOR.MIDDLE)

            pry, prh = band(0.32, 0.52)
            price_pt = fit_size(price, inner, 15, 10.5, 1, HEAD)
            text(slide, x + pad_x, pry, inner, prh,
                 [[(price, HEAD, price_pt, INK, False)]],
                 align=PP_ALIGN.CENTER, space_after=0, anchor=MSO_ANCHOR.MIDDLE)

            pky, pkh = band(0.52, 0.62)
            text(slide, x + pad_x, pky, inner, pkh,
                 [[(pack, BODY, 7, SAND_MED, False)]],
                 align=PP_ALIGN.CENTER, space_after=0, anchor=MSO_ANCHOR.MIDDLE)

            bgy, bgh = band(0.64, 0.78)
            badge_w = min(inner * 0.92, max(0.90, 0.09 * max(len(badge), 1) + 0.40))
            badge_x = x + (card_w - badge_w) / 2
            badge_h = min(0.20, bgh)
            rrect(slide, badge_x, bgy + (bgh - badge_h) / 2, badge_w, badge_h,
                  FOREST, None, radius=0.40)
            text(slide, badge_x, bgy + (bgh - badge_h) / 2, badge_w, badge_h,
                 [[(badge, BODY, 7.5, WHITE, True)]],
                 align=PP_ALIGN.CENTER, space_after=0, anchor=MSO_ANCHOR.MIDDLE)

            if feature:
                fy, fh = band(0.80, 0.97)
                feat_pt = fit_size(feature, inner, 7.5, 6, 2)
                text(slide, x + pad_x, fy, inner, fh,
                     [[(feature, BODY, feat_pt, SAND_DARK, False)]],
                     align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.05,
                     anchor=MSO_ANCHOR.MIDDLE)

    footer(slide, ctx.deck_name, ctx.page)


def r_fullimage(slide, spec: Slide, ctx: Ctx):
    """Helbreddsbild under rubrik — för infografik och jämförelser."""
    background(slide)
    y = title_block(slide, spec, base_pt=24, gap=0.04)
    img = image_path(spec.images[0]) if spec.images else None
    bottom = FOOTER_Y - (0.48 if spec.caption else 0.28)
    if img:
        image_contain(slide, img, MARGIN - 0.10, y, CONTENT_W + 0.20, bottom - y)
    if spec.caption:
        caption(slide, spec, FOOTER_Y - 0.48)
    footer(slide, ctx.deck_name, ctx.page)


def r_close(slide, spec: Slide, ctx: Ctx):
    background(slide, INK)
    rect(slide, 0, 0, 0.10, SH, FOREST)

    logo = os.path.join(PUBLIC, "brand", "roots-logo-white.png")
    if os.path.exists(logo):
        slide.shapes.add_picture(logo, Inches(SW - MARGIN - 1.35), Inches(0.78),
                                 height=Inches(0.38))

    kicker(slide, spec.kicker, x=1.30, y=1.30, color=OLIVE)

    width = 9.2
    size = fit_size(spec.title, width, 38, 22, 2, HEAD)
    lines = line_count(spec.title, width, size, HEAD)
    title_h = lines * size * 1.16 / 72
    text(slide, 1.30, 1.72, width, title_h + 0.1,
         [[(spec.title, HEAD, size, WHITE, False)]], space_after=0, line_spacing=1.12)

    y = 1.72 + title_h + 0.30
    hairline(slide, 1.30, y, 1.5, OLIVE, 0.022)
    bullets(slide, spec.items, 1.30, y + 0.44, 9.6, base_pt=14.5, gap_extra=0.22,
            dot_color=OLIVE, color=SAND_LIGHT, max_h=FOOTER_Y - 0.5 - y)

    footer(slide, ctx.deck_name, ctx.page, light=True)


def r_bignumber(slide, spec: Slide, ctx: Ctx):
    """En dominerande siffra med förklaring och tre roller under.

    Master Source MS-003: "En stor siffra som blickfång, med tre mindre
    produktroller nedanför."
    """
    background(slide)
    y = title_block(slide, spec)

    figure = spec.items[0] if spec.items else ""
    explain = spec.items[1] if len(spec.items) > 1 else ""
    roles = [tuple(p.strip() for p in item.split("|")) + ("",)
             for item in spec.items[2:]]

    fig_pt = fit_size(figure, CONTENT_W, 66, 34, 1, HEAD)
    text(slide, MARGIN, y + 0.10, CONTENT_W, fig_pt / 72 * 1.35,
         [[(figure, HEAD, fig_pt, FOREST, False)]], space_after=0)

    top = y + 0.10 + fig_pt / 72 * 1.28
    if explain:
        # Bredare förklaringsrad så sista orden inte hänger ensamma.
        exp_w = min(10.2, CONTENT_W)
        exp_pt = fit_size(explain, exp_w, 15, 11.5, 3)
        text(slide, MARGIN, top, exp_w, 0.85,
             [[(explain, BODY, exp_pt, SAND_DARK, False)]], space_after=0,
             line_spacing=1.32)
        top += line_count(explain, exp_w, exp_pt) * exp_pt * 1.32 / 72 + 0.34

    n = max(1, len(roles))
    gap = 0.30
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = min(1.45, FOOTER_Y - 0.40 - top)
    for i, role in enumerate(roles):
        name, note = role[0], role[1]
        x = MARGIN + i * (card_w + gap)
        rrect(slide, x, top, card_w, card_h, SAND_100)
        name_pt = fit_size(name, card_w - 0.6, 16, 12, 1, HEAD)
        name_band = name_pt / 72 * 1.55 + 0.10
        text(slide, x + 0.30, top + 0.22, card_w - 0.6, name_band,
             [[(name, HEAD, name_pt, INK, False)]], space_after=0,
             anchor=MSO_ANCHOR.MIDDLE)
        note_top = top + 0.22 + name_band + 0.10
        note_pt = fit_size(note, card_w - 0.6, 12.5, 9.5, 2)
        text(slide, x + 0.30, note_top, card_w - 0.6,
             card_h - (note_top - top) - 0.22,
             [[(note, BODY, note_pt, SAND_DARK, False)]],
             space_after=0, line_spacing=1.3)

    footer(slide, ctx.deck_name, ctx.page)
    caption(slide, spec, FOOTER_Y - 0.42)


def r_words(slide, spec: Slide, ctx: Ctx):
    """Några få ord i stor grad, ett per fält.

    För slides där Master Source anger byggstenar utan beskrivningar —
    MS-009 "Tre premiumpelare", MS-011 fyra informationskort.
    """
    background(slide)
    y = title_block(slide, spec)

    words = [w.strip() for w in spec.items if w.strip()]
    n = max(1, len(words))
    gap = 0.34
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = min(3.10, FOOTER_Y - 0.75 - y)

    for i, word in enumerate(words):
        x = MARGIN + i * (card_w + gap)
        rrect(slide, x, y, card_w, card_h, WHITE, SAND_LIGHT)
        rect(slide, x, y, card_w, 0.035, FOREST)
        size = fit_size(word, card_w - 0.6, 26, 14, 2, HEAD)
        text(slide, x + 0.30, y, card_w - 0.6, card_h,
             [[(word, HEAD, size, INK, False)]], align=PP_ALIGN.CENTER,
             space_after=0, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.16)

    footer(slide, ctx.deck_name, ctx.page)
    caption(slide, spec, y + card_h + 0.26)


def r_calc(slide, spec: Slide, ctx: Ctx):
    """Kalkylpanel till vänster, resultatkort till höger.

    Master Source MS-015: "Delad layout: Vänster: enkel kalkylpanel. Höger:
    resultatkort med tydlig ekonomisk potential."
    """
    background(slide)
    y = title_block(slide, spec)

    gap = 0.42
    left_w = (CONTENT_W - gap) * 0.56
    right_w = CONTENT_W - gap - left_w
    panel_h = min(3.55, FOOTER_Y - 0.40 - y)

    rrect(slide, MARGIN, y, left_w, panel_h, WHITE, SAND_LIGHT)
    # items: "etikett" eller "etikett | belopp" — belopp visas i högerfältet.
    fields = []
    for raw in spec.items:
        if not raw.strip():
            continue
        parts = [p.strip() for p in raw.split("|", 1)]
        fields.append((parts[0], parts[1] if len(parts) > 1 else ""))
    row_h = (panel_h - 0.70) / max(1, len(fields))
    for i, (name, amount) in enumerate(fields):
        top = y + 0.35 + i * row_h
        label_pt = fit_size(name, left_w - 2.05, 13.5, 10.5, 2)
        text(slide, MARGIN + 0.40, top, left_w - 2.05, row_h,
             [[(name, BODY, label_pt, INK, False)]], space_after=0,
             anchor=MSO_ANCHOR.MIDDLE)
        box_x = MARGIN + left_w - 1.55
        rrect(slide, box_x, top + row_h / 2 - 0.18, 1.20, 0.36, SAND_100)
        if amount:
            amt_pt = fit_size(amount, 1.05, 11.5, 9, 1)
            text(slide, box_x + 0.06, top + row_h / 2 - 0.16, 1.08, 0.32,
                 [[(amount, BODY, amt_pt, FOREST, True)]], space_after=0,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Inga hairlines mellan rader — beloppsrutor + linjer blir stökigt.

    rx = MARGIN + left_w + gap
    rrect(slide, rx, y, right_w, panel_h, FOREST_SOFT)
    rect(slide, rx, y, right_w, 0.04, FOREST)
    head = spec.columns[0] if spec.columns else ""
    body = spec.columns[1] if len(spec.columns) > 1 else ""
    head_pt = fit_size(head, right_w - 0.9, 20, 14, 3, HEAD)
    head_h = line_count(head, right_w - 0.9, head_pt, HEAD) * head_pt * 1.2 / 72
    text(slide, rx + 0.45, y + 0.50, right_w - 0.9, head_h + 0.1,
         [[(head, HEAD, head_pt, INK, False)]], space_after=0, line_spacing=1.18)
    body_pt = fit_size(body, right_w - 0.9, 13.5, 10.5, 6)
    text(slide, rx + 0.45, y + 0.58 + head_h, right_w - 0.9,
         panel_h - head_h - 1.0,
         [[(body, BODY, body_pt, SAND_DARK, False)]], space_after=0,
         line_spacing=1.34)

    footer(slide, ctx.deck_name, ctx.page)


def r_logo(slide, spec: Slide, ctx: Ctx):
    """Rent märkesomslag: logotypen ensam, centrerad på mörk botten."""
    background(slide, INK)
    rect(slide, 0, 0, 0.10, SH, FOREST)

    logo = os.path.join(PUBLIC, "brand", "roots-logo-white.png")
    if os.path.exists(logo):
        pic = slide.shapes.add_picture(logo, Inches(0), Inches(0),
                                       height=Inches(1.15))
        pic.left = Inches(SW / 2) - pic.width // 2
        pic.top = Inches(SH / 2 - 1.05)

    if spec.kicker:
        kicker(slide, spec.kicker, x=0, y=SH / 2 - 1.72, color=OLIVE, w=SW,
               align=PP_ALIGN.CENTER)

    y = SH / 2 + 0.44
    hairline(slide, SW / 2 - 0.75, y, 1.5, FOREST, 0.022)

    if spec.title:
        size = fit_size(spec.title, CONTENT_W, 22, 15, 2, HEAD)
        text(slide, MARGIN, y + 0.40, CONTENT_W, 0.8,
             [[(spec.title, HEAD, size, WHITE, False)]], align=PP_ALIGN.CENTER,
             space_after=0, line_spacing=1.2)
    if spec.subtitle:
        text(slide, MARGIN, y + 1.10, CONTENT_W, 0.6,
             [[(spec.subtitle, BODY, 13, SAND_MED, False)]],
             align=PP_ALIGN.CENTER, space_after=0)


def r_story(slide, spec: Slide, ctx: Ctx):
    """Berättarslide: punkter till vänster, ett bildfält till höger.

    Bilden läggs hel i sitt fält i stället för beskuren, eftersom bilderna i
    berättelsen är föremål och porträtt som inte tål att kapas i kanterna.
    """
    background(slide)
    panel_w, panel_x = 4.30, SW - MARGIN - 4.30
    body_w = panel_x - MARGIN - 0.55

    y = title_block(slide, spec, width=body_w)

    img = image_path(spec.images[0]) if spec.images else None
    if img:
        panel_h = 3.30
        panel_y = (SH - panel_h) / 2
        rrect(slide, panel_x, panel_y, panel_w, panel_h, WHITE, SAND_LIGHT)
        image_contain(slide, img, panel_x + 0.30, panel_y + 0.30,
                      panel_w - 0.60, panel_h - 0.60)

    cap_room = 0.98 if spec.caption else 0
    bullets(slide, spec.items, MARGIN, y, body_w, base_pt=14,
            max_h=FOOTER_Y - 0.35 - y - cap_room)
    caption(slide, spec, FOOTER_Y - 1.02)
    footer(slide, ctx.deck_name, ctx.page)


def r_pitch(slide, spec: Slide, ctx: Ctx):
    """Produktslide: bild till höger, två fält till vänster.

    Fälten är säljarens två verktyg — det som är sant om produkten och det
    kunden känner. De delar höjden efter hur många punkter de innehåller.
    """
    background(slide)
    img = image_path(spec.images[0]) if spec.images else None
    img_w, img_x = 3.85, SW - MARGIN - 3.85
    body_w = img_x - 0.55 - MARGIN if img else CONTENT_W
    bottom = FOOTER_Y - 0.30

    y = title_block(slide, spec, width=body_w)
    if img:
        # Säljfrasen ligger under bilden i stället för under punkterna: den
        # vänstra spalten behöver all höjd den kan få.
        img_h = bottom - 1.24 - (1.28 if spec.caption else 0)
        image_cover(slide, img, img_x, 1.24, img_w, img_h, rounded=True,
                    focus_y=0.40)
        if spec.caption:
            top = 1.24 + img_h + 0.24
            rrect(slide, img_x, top, img_w, bottom - top, FOREST_SOFT)
            size = fit_size(spec.caption, img_w - 0.56, 12.5, 9.5, 5)
            text(slide, img_x + 0.28, top + 0.20, img_w - 0.56,
                 bottom - top - 0.32,
                 [[(spec.caption, BODY, size, FOREST, False)]], space_after=0,
                 line_spacing=1.3)

    blocks = [c for c in (spec.left, spec.right) if c]
    if not blocks:
        footer(slide, ctx.deck_name, ctx.page)
        return

    gap = 0.20
    # Radhöjden i bullets är en uppskattning, så nedre marginalen är tilltagen:
    # utan den lägger sig sista punkten på kortets underkant.
    head_room, pad = 0.52, 0.28
    space = bottom - y - gap * (len(blocks) - 1)
    # Rubriken kostar lika mycket i ett kort fält som i ett långt. Utan den i
    # vikten blir det korta fältet så pressat att punkterna hamnar utanför.
    weights = [len(b.items) + 1.5 for b in blocks]

    top = y
    for block, weight in zip(blocks, weights):
        h = space * weight / sum(weights)
        rrect(slide, MARGIN, top, body_w, h, WHITE, SAND_LIGHT)
        rect(slide, MARGIN, top, body_w, 0.032, FOREST)
        kicker(slide, block.heading, x=MARGIN + 0.34, y=top + 0.22,
               color=FOREST, w=body_w - 0.68)
        bullets(slide, block.items, MARGIN + 0.34, top + head_room,
                body_w - 0.68, base_pt=12.5, gap_extra=0.12,
                max_h=h - head_room - pad)
        top += h + gap

    if not img:
        caption(slide, spec, bottom + 0.18, w=body_w)
    footer(slide, ctx.deck_name, ctx.page)


def r_skarm(slide, spec: Slide, ctx: Ctx):
    """Skärmbild ur plattformen till höger, stödord till vänster.

    Bilden får styra sin egen bredd: en portalvy är bred och låg, en telefon
    smal och hög. Telefonvyerna lämnar därför nästan dubbelt så mycket plats
    åt texten, vilket är rimligt — de har färre saker att peka på.
    """
    background(slide)

    rel = spec.images[0] if spec.images else None
    src = os.path.join(ASSETS, "skarm", rel) if rel else None
    portrait = False
    if src and os.path.exists(src):
        iw, ih = Image.open(src).size
        portrait = ih > iw

    # Sidomenyn i portalen går ända ned; hoppa över den när vi letar efter
    # var innehållet slutar.
    img = screenshot_path(rel, 0.0 if portrait else 0.22) if rel else None

    max_w = 2.62 if portrait else 7.00
    pad = 0.10
    top_limit, bottom_limit = 1.02, 6.58
    max_h = bottom_limit - top_limit - 2 * pad

    draw_w = draw_h = 0.0
    if img:
        iw, ih = Image.open(img).size
        draw_w = min(max_w, max_h * iw / ih)
        draw_h = draw_w * ih / iw

    panel_w = draw_w + 2 * pad
    panel_x = SW - MARGIN - panel_w
    column_w = (panel_x - MARGIN - 0.60) if img else CONTENT_W
    # En telefonbild lämnar över åtta tum till texten, och en punktlista som
    # löper så brett blir en rad tunna streck. Radlängden får därför ett tak
    # även när spalten är bred.
    body_w = min(column_w, 5.90)

    # Rubriken får däremot bruka hela spalten — den tål långa rader.
    y = title_block(slide, spec, width=column_w, base_pt=31, gap=0.06)

    if img:
        panel_h = draw_h + 2 * pad
        panel_y = top_limit + (bottom_limit - top_limit - panel_h) / 2
        rrect(slide, panel_x, panel_y, panel_w, panel_h, WHITE, SAND_LIGHT,
              radius=0.03)
        slide.shapes.add_picture(img, Inches(panel_x + pad),
                                 Inches(panel_y + pad), Inches(draw_w),
                                 Inches(draw_h))

    cap_room = 0.92 if spec.caption else 0
    floor = FOOTER_Y - 0.35 - cap_room
    _, list_h = bullets_extent(spec.items, body_w, 14, 0.20, floor - y)
    # Lägg listan i höjd med bildens tyngdpunkt i stället för att låta den
    # hänga kvar under rubriken med tom yta under sig.
    top = y + max(0.0, (floor - y - list_h) * 0.42)
    bullets(slide, spec.items, MARGIN, top, body_w, base_pt=14,
            max_h=floor - y)
    if spec.caption:
        caption(slide, spec, FOOTER_Y - 0.98, w=body_w)
    footer(slide, ctx.deck_name, ctx.page)


def r_phones(slide, spec: Slide, ctx: Ctx):
    """2–3 mobilskärmar sida vid sida — en roll, flera steg i fickan.

    `images` är filnamn under `_assets/skarm/` (porträtt). `items` är etiketter
    under respektive telefon. Färre bilder ger större telefoner.
    """
    background(slide)
    y = title_block(slide, spec, base_pt=28, gap=0.02)

    rels = list(spec.images[:3])
    labels = list(spec.items[: len(rels)])
    while len(labels) < len(rels):
        labels.append("")

    n = max(1, len(rels))
    gap = 0.50 if n <= 2 else 0.36
    col_w = (CONTENT_W - gap * (n - 1)) / n
    floor = FOOTER_Y - (0.68 if spec.caption else 0.28)
    label_h = 0.34 if any(labels) else 0.0
    frame_h = min(5.05 if n <= 2 else 4.75, floor - y - label_h - 0.08)

    bezel = 0.10
    for i, rel in enumerate(rels):
        x = MARGIN + i * (col_w + gap)
        img = screenshot_path(rel, 0.0)
        # Två telefoner får vara bredare; tre håller klassisk proportion.
        phone_w = min(col_w * 0.92, frame_h * (0.58 if n <= 2 else 0.50))
        phone_h = frame_h
        px = x + (col_w - phone_w) / 2

        rrect(slide, px, y, phone_w, phone_h, INK, None, radius=0.12)
        sx = px + bezel
        sy = y + bezel + 0.07
        sw = phone_w - 2 * bezel
        sh = phone_h - 2 * bezel - 0.14
        rrect(slide, sx, sy, sw, sh, WHITE, None, radius=0.055)

        if img:
            iw, ih = Image.open(img).size
            draw_w = sw
            draw_h = draw_w * ih / iw
            if draw_h > sh:
                draw_h = sh
                draw_w = draw_h * iw / ih
            pic_x = sx + (sw - draw_w) / 2
            pic_y = sy + (sh - draw_h) / 2
            slide.shapes.add_picture(img, Inches(pic_x), Inches(pic_y),
                                     Inches(draw_w), Inches(draw_h))

        rect(slide, px + phone_w / 2 - 0.26, y + phone_h - 0.13, 0.52, 0.032,
             RGBColor(0x3A, 0x3A, 0x38))

        if labels[i]:
            lab_pt = fit_size(labels[i], col_w, 13.5, 11, 2, HEAD)
            text(slide, x, y + phone_h + 0.12, col_w, label_h,
                 [[(labels[i], HEAD, lab_pt, INK, False)]],
                 align=PP_ALIGN.CENTER, space_after=0)

    if spec.caption:
        caption(slide, spec, FOOTER_Y - 0.68)
    footer(slide, ctx.deck_name, ctx.page)


def r_desktops(slide, spec: Slide, ctx: Ctx):
    """Två skrivbordsskärmar sida vid sida — förening / lagledare.

    `images` är landskaps-PNG under `_assets/skarm/`. `items` etiketter under.
    """
    background(slide)
    y = title_block(slide, spec, base_pt=28, gap=0.04)

    rels = list(spec.images[:2])
    labels = list(spec.items[: len(rels)])
    while len(labels) < len(rels):
        labels.append("")

    n = max(1, len(rels))
    gap = 0.34
    col_w = (CONTENT_W - gap * (n - 1)) / n
    floor = FOOTER_Y - (0.78 if spec.caption else 0.28)
    label_h = 0.32 if any(labels) else 0.0
    panel_h = min(4.35, floor - y - label_h - 0.10)
    pad = 0.10

    for i, rel in enumerate(rels):
        x = MARGIN + i * (col_w + gap)
        img = screenshot_path(rel, 0.22)
        rrect(slide, x, y, col_w, panel_h, WHITE, SAND_LIGHT, radius=0.03)
        # Tunn “browser chrome”
        rect(slide, x, y, col_w, 0.22, SAND_100)
        oval(slide, x + 0.14, y + 0.07, 0.08, 0.08, SAND_MED)
        oval(slide, x + 0.28, y + 0.07, 0.08, 0.08, SAND_MED)
        oval(slide, x + 0.42, y + 0.07, 0.08, 0.08, SAND_LIGHT)

        inner_w = col_w - 2 * pad
        inner_h = panel_h - 0.22 - 2 * pad
        if img:
            iw, ih = Image.open(img).size
            draw_w = inner_w
            draw_h = draw_w * ih / iw
            if draw_h > inner_h:
                draw_h = inner_h
                draw_w = draw_h * iw / ih
            pic_x = x + pad + (inner_w - draw_w) / 2
            pic_y = y + 0.22 + pad + (inner_h - draw_h) / 2
            slide.shapes.add_picture(img, Inches(pic_x), Inches(pic_y),
                                     Inches(draw_w), Inches(draw_h))

        if labels[i]:
            lab_pt = fit_size(labels[i], col_w, 13, 10.5, 2, HEAD)
            text(slide, x, y + panel_h + 0.12, col_w, label_h,
                 [[(labels[i], HEAD, lab_pt, INK, False)]],
                 align=PP_ALIGN.CENTER, space_after=0)

    if spec.caption:
        caption(slide, spec, FOOTER_Y - 0.68)
    footer(slide, ctx.deck_name, ctx.page)


def r_impact(slide, spec: Slide, ctx: Ctx):
    """Maffig säljsiffra på mörk botten — en siffra som bär hela sliden.

    `title` är den stora siffran (t.ex. »35 %«), `subtitle` raden under,
    `items` tre korta påståenden i rad, `caption` avslutningsrad.
    """
    background(slide, INK)
    rect(slide, 0, 0, 0.12, SH, FOREST)

    # Mjuk sandyta som ger djup utan gradient-prål.
    rect(slide, SW - 4.2, 0, 4.2, SH, RGBColor(0x24, 0x24, 0x21))

    kicker(slide, spec.kicker, x=MARGIN, y=0.78, color=OLIVE, w=CONTENT_W)

    fig = spec.title or ""
    fig_pt = fit_size(fig, CONTENT_W - 0.4, 128, 72, 1, HEAD)
    fig_h = fig_pt * 1.05 / 72
    text(slide, MARGIN, 1.55, CONTENT_W - 0.4, fig_h + 0.15,
         [[(fig, HEAD, fig_pt, WHITE, False)]], space_after=0, line_spacing=0.95)

    y = 1.55 + fig_h + 0.10
    if spec.subtitle:
        sub_pt = fit_size(spec.subtitle, 9.5, 28, 18, 2, HEAD)
        sub_h = line_count(spec.subtitle, 9.5, sub_pt, HEAD) * sub_pt * 1.15 / 72
        text(slide, MARGIN, y, 9.5, sub_h + 0.1,
             [[(spec.subtitle, HEAD, sub_pt, OLIVE, False)]], space_after=0,
             line_spacing=1.12)
        y += sub_h + 0.28

    hairline(slide, MARGIN, y, 1.8, FOREST, 0.028)
    y += 0.42

    points = [p.strip() for p in spec.items if p.strip()][:3]
    if points:
        gap = 0.28
        col_w = (CONTENT_W - gap * (len(points) - 1)) / len(points)
        for i, point in enumerate(points):
            x = MARGIN + i * (col_w + gap)
            # Numerisk chip
            oval(slide, x, y + 0.06, 0.28, 0.28, FOREST)
            text(slide, x, y + 0.08, 0.28, 0.24,
                 [[(f"{i + 1}", BODY, 10, WHITE, True)]],
                 align=PP_ALIGN.CENTER, space_after=0)
            body_pt = fit_size(point, col_w - 0.10, 14.5, 11, 3)
            text(slide, x, y + 0.48, col_w, 1.35,
                 [[(point, BODY, body_pt, SAND_LIGHT, False)]],
                 space_after=0, line_spacing=1.32)

    if spec.caption:
        text(slide, MARGIN, FOOTER_Y - 0.55, CONTENT_W, 0.35,
             [[(spec.caption, BODY, 12.5, OLIVE, False)]], space_after=0)

    footer(slide, ctx.deck_name, ctx.page, light=True)


def r_team(slide, spec: Slide, ctx: Ctx):
    """Teamgrid: enhetliga porträtt med namn under.

    Minimalistiskt — bilden kant i kant, tunn forest-accent, ingen vit
    passpartout. `images` och `items` är parallella listor.
    """
    background(slide)
    # Kompakt rubrik — varje tiondels tum går till porträtten.
    y = title_block(slide, spec, base_pt=25, gap=0.0)

    names = list(spec.items)
    paths = [image_path(p) for p in spec.images]
    people = [(n, p) for n, p in zip(names, paths) if n]
    n = max(1, len(people))
    # Sju personer: 4 + 3. Annars max tre per rad så ansiktena håller storlek.
    if n >= 7:
        per_row = 4
    elif n > 4:
        per_row = 3
    else:
        per_row = n
    row_count = (n + per_row - 1) // per_row

    # Lite smalare sidmarginal än övriga slides så ansiktena blir större.
    side = 0.55
    grid_w = SW - 2 * side
    gap_x, gap_y = 0.16, 0.08
    name_h = 0.24
    # Caption behöver mer plats än en tunn rad — annars lägger den sig
    # över nedersta porträttraden.
    cap_h = 0.72 if spec.caption else 0
    bottom = FOOTER_Y - 0.18 - cap_h
    avail_h = bottom - y - gap_y * (row_count - 1)

    tile_w = (grid_w - gap_x * (per_row - 1)) / per_row
    # Väx till den storlek som fyller antingen bredd eller höjd (4:5).
    photo_h_cap = avail_h / row_count - name_h
    photo_w_from_h = photo_h_cap * 4 / 5
    if photo_w_from_h <= tile_w and photo_h_cap > 1.0:
        photo_w, photo_h = photo_w_from_h, photo_h_cap
    else:
        photo_w = tile_w
        photo_h = photo_w * 5 / 4
        if photo_h + name_h > avail_h / row_count:
            photo_h = avail_h / row_count - name_h
            photo_w = photo_h * 4 / 5

    for i, (name, path) in enumerate(people):
        col, row = i % per_row, i // per_row
        in_row = min(per_row, n - row * per_row)
        row_w = in_row * photo_w + gap_x * (in_row - 1)
        row_x0 = side + (grid_w - row_w) / 2
        x = row_x0 + col * (photo_w + gap_x)
        top = y + row * (photo_h + name_h + gap_y)

        # Bilden fyller rutan — ingen vit ram/padding.
        if path:
            image_cover(slide, path, x, top, photo_w, photo_h,
                        rounded=True, focus_y=0.28)
        else:
            rrect(slide, x, top, photo_w, photo_h, SAND_100, None, radius=0.03)
        # Endast en tunn forest-linje som accent, inte en hel ram.
        rect(slide, x, top, photo_w, 0.016, FOREST)

        name_pt = fit_size(name, photo_w - 0.06, 12, 9.5, 1, HEAD)
        text(slide, x, top + photo_h + 0.04, photo_w, name_h,
             [[(name, HEAD, name_pt, INK, False)]],
             align=PP_ALIGN.CENTER, space_after=0)

    if spec.caption:
        caption(slide, spec, FOOTER_Y - 0.70)
    footer(slide, ctx.deck_name, ctx.page)


RENDERERS = {
    "cover": r_cover,
    "logo": r_logo,
    "story": r_story,
    "pitch": r_pitch,
    "skarm": r_skarm,
    "phones": r_phones,
    "desktops": r_desktops,
    "impact": r_impact,
    "bignumber": r_bignumber,
    "words": r_words,
    "calc": r_calc,
    "hero": r_hero,
    "points": r_points,
    "cards": r_cards,
    "columns": r_columns,
    "timeline": r_timeline,
    "quarters": r_quarters,
    "years": r_years,
    "flywheel": r_flywheel,
    "products": r_products,
    "table": r_table,
    "kpi": r_kpi,
    "roles": r_roles,
    "team": r_team,
    "fullimage": r_fullimage,
    "compare": r_compare,
    "close": r_close,
}
