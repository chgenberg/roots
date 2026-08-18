#!/usr/bin/env python3
"""Bygger Roots tre presentationer ur Master Source.

    python3 docs/presentations/build_decks.py            # alla tre
    python3 docs/presentations/build_decks.py d1         # bara säljpresentationen

Innehållet kommer ur kollegans slide-specifikationer via decks.py — ingen text
skrivs om på vägen — och designen ur roots_deck.py. Talarmanus och källhänvisning
läggs i respektive slides anteckningsfält, så att en säljare kan skriva ut
presentationen med manus.

Kör verify_verbatim.py efteråt för att kontrollera att varje textrad på sliderna
finns ordagrant i källdokumenten.
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from pptx.util import Inches  # noqa: E402

import roots_deck as rd  # noqa: E402
from decks import DECKS as MS_DECKS  # noqa: E402
import deck_salj  # noqa: E402
import deck_forening  # noqa: E402
import deck_kommunikation  # noqa: E402
from slide_model import Deck  # noqa: E402

HERE = Path(__file__).parent

# Sidfotens text per presentation.
FOOTERS = {
    "Roots_Saljpresentation.pptx": "Roots · Säljpresentation",
    "Roots_Roadmap_Ar_1-5.pptx": "Roots · Roadmap år 1–5",
    "Roots_Oversiktspresentation.pptx": "Roots · Det här är Roots",
    "Roots_Ar_1_Execution_Roadmap.pptx": "Roots · År 1 – Execution Roadmap",
    "Roots_Bankpresentation.pptx": "Roots · Bankunderlag",
    "Roots_Foreningspresentation.pptx": "Roots · Till föreningen",
    "Roots_Kommunikationskit.pptx": "Roots · Kommunikationskit",
    "Roots_Plattformen.pptx": "Roots · Plattformen",
}
KEYS = {"d1": "Roots_Saljpresentation.pptx",
        "d2": "Roots_Roadmap_Ar_1-5.pptx",
        "d3": "Roots_Oversiktspresentation.pptx",
        "d4": "Roots_Ar_1_Execution_Roadmap.pptx",
        "forening": "Roots_Foreningspresentation.pptx",
        "komm": "Roots_Kommunikationskit.pptx"}

# D1 kommer ur deck_salj (noter + MS). Förenings-/komm-decks är fristående.
DECKS = {
    **MS_DECKS,
    "Roots_Saljpresentation.pptx": deck_salj.build,
    "Roots_Foreningspresentation.pptx": deck_forening.build,
    "Roots_Kommunikationskit.pptx": deck_kommunikation.build,
}


def register_notes_master(prs: Presentation) -> None:
    """Skriver in anteckningsmallen i presentation.xml.

    Så fort en slide får talarmanus skapar python-pptx (1.0.2) en notesMaster-del
    och en relation till den, men lägger aldrig till <p:notesMasterIdLst>.
    PowerPoint och LibreOffice struntar i det; Keynote avvisar hela filen med
    "Filformatet är ogiltigt". Elementet måste ligga direkt efter
    sldMasterIdLst enligt schemat.
    """
    rid = next((rel.rId for rel in prs.part.rels.values()
                if rel.reltype.endswith("/notesMaster")), None)
    root = prs._element
    if rid is None or root.find(qn("p:notesMasterIdLst")) is not None:
        return

    lst = root.makeelement(qn("p:notesMasterIdLst"), {})
    lst.append(root.makeelement(qn("p:notesMasterId"), {qn("r:id"): rid}))
    root.find(qn("p:sldMasterIdLst")).addnext(lst)


def render(deck: Deck, out: Path, deck_name: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(rd.SW)
    prs.slide_height = Inches(rd.SH)
    blank = prs.slide_layouts[6]

    for index, spec in enumerate(deck.slides, start=1):
        renderer = rd.RENDERERS.get(spec.layout)
        if renderer is None:
            raise ValueError(
                f"{out.name}: okänd layout {spec.layout!r} på slide {spec.number}")

        slide = prs.slides.add_slide(blank)
        renderer(slide, spec, rd.Ctx(deck_name=deck_name, page=index))

        if spec.notes:
            slide.notes_slide.notes_text_frame.text = spec.notes

    register_notes_master(prs)
    prs.save(out)
    size_mb = out.stat().st_size / 1_048_576
    print(f"  {len(deck.slides):>2} slides  {size_mb:>5.1f} MB  {out.name}")


def main() -> None:
    wanted = {KEYS.get(a.lower(), a) for a in sys.argv[1:]}
    print("Bygger presentationer ur Master Source (docs/knowledge-os/_extract)\n")

    for filename, build in DECKS.items():
        if wanted and filename not in wanted:
            continue
        render(build(), HERE / filename, FOOTERS[filename])

    print(f"\nKlart. Filerna ligger i {HERE}")


if __name__ == "__main__":
    main()
