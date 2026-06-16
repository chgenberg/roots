#!/usr/bin/env python3
"""Bygger Roots pitch-deck (.pptx) för föreningar/kund."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))
WEB = os.path.abspath(os.path.join(ROOT, "..", "..", "apps", "web", "public"))
ASSETS = os.path.join(ROOT, "assets")
OUT = os.path.join(ROOT, "Roots_Pitch_Deck.pptx")

INK = RGBColor(0x1D, 0x1D, 0x1B)
SAND_DARK = RGBColor(0x7F, 0x71, 0x5B)
SAND_MED = RGBColor(0xB2, 0xA4, 0x91)
SAND_LIGHT = RGBColor(0xD5, 0xCA, 0xBF)
SAND_100 = RGBColor(0xF1, 0xEB, 0xE2)
OFFWHITE = RGBColor(0xFA, 0xF6, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOREST = RGBColor(0x6B, 0x79, 0x4F)
OLIVE = RGBColor(0xC1, 0xBF, 0x99)
SKY = RGBColor(0xA7, 0xBB, 0xC5)
TERRA = RGBColor(0xE1, 0x87, 0x54)

HEAD = "Alan Sans"
BODY = "Inter 18pt"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=OFFWHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def rrect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = 0.06
    except Exception:
        pass
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            text, font, size, color, bold = (seg + (False,))[:5] if len(seg) < 5 else seg
            r = p.add_run(); r.text = text
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold
    return tb


def kicker(slide, x, y, text, color=SAND_DARK, w=Inches(8)):
    return txt(slide, x, y, w, Inches(0.4),
               [[(text.upper(), BODY, 12, color, True)]], space_after=0)


def pic(slide, name, x, y, w=None, h=None, from_web=False):
    path = os.path.join(WEB if from_web else ASSETS, name)
    if not os.path.exists(path):
        return None
    kwargs = {}
    if w: kwargs["width"] = w
    if h: kwargs["height"] = h
    return slide.shapes.add_picture(path, x, y, **kwargs)


def pic_fit(slide, name, x, y, max_w, max_h, from_web=False, align="center", valign="middle"):
    """Place image scaled to fit within a box, centered."""
    path = os.path.join(WEB if from_web else ASSETS, name)
    if not os.path.exists(path):
        return None
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = int(max_w / ar)
    else:
        h = max_h; w = int(max_h * ar)
    if align == "center":
        px = x + (max_w - w) // 2
    elif align == "left":
        px = x
    else:
        px = x + (max_w - w)
    if valign == "middle":
        py = y + (max_h - h) // 2
    elif valign == "top":
        py = y
    else:
        py = y + (max_h - h)
    return slide.shapes.add_picture(path, px, py, width=w, height=h)


def dot(slide, x, y, d, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def page_no(slide, n):
    txt(slide, SW - Inches(0.9), SH - Inches(0.55), Inches(0.6), Inches(0.35),
        [[(str(n), BODY, 10, SAND_MED, False)]], align=PP_ALIGN.RIGHT, space_after=0)


def footer(slide):
    txt(slide, Inches(0.6), SH - Inches(0.55), Inches(8), Inches(0.35),
        [[("ROOTS", HEAD, 11, SAND_DARK, True),
          ("   \u00b7   Riktigt bra produkter. Starkare föreningsliv.", BODY, 10, SAND_MED, False)]],
        space_after=0)


def new_slide(color=OFFWHITE):
    s = prs.slides.add_slide(BLANK)
    bg(s, color)
    return s


def section_head(slide, kick, title, title_color=INK, y=Inches(0.7), tw=Inches(11.8)):
    kicker(slide, Inches(0.7), y, kick)
    txt(slide, Inches(0.66), y + Inches(0.42), tw, Inches(1.4),
        [[(title, HEAD, 30, title_color, True)]], space_after=0, line_spacing=1.05)


# ───────────────────────── 1. COVER ─────────────────────────
def slide_cover():
    s = new_slide(OFFWHITE)
    rect(s, 0, 0, Inches(0.28), SH, FOREST)
    pic_fit(s, "brand/roots-logo-black.png", Inches(0.9), Inches(0.8),
            Inches(2.6), Inches(1.0), from_web=True, align="left")
    txt(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.6),
        [[("Samla in pengar med", HEAD, 40, INK, True)],
         [("produkter folk faktiskt vill ha ", HEAD, 40, INK, True), ("igen.", HEAD, 40, FOREST, True)]],
        space_after=4, line_spacing=1.04)
    txt(s, Inches(0.95), Inches(4.9), Inches(10.5), Inches(1.2),
        [[("Nordisk hår- och hudvård i premiumklass \u2013 forskningsförankrad, snäll mot håret, "
           "och gjord för att stärka svenskt föreningsliv.", BODY, 16, SAND_DARK, False)]],
        line_spacing=1.3)
    txt(s, Inches(0.95), Inches(6.5), Inches(10), Inches(0.5),
        [[("Ett förslag till er förening", BODY, 12, SAND_MED, False)]], space_after=0)


# ───────────────────────── 2. PROBLEM ─────────────────────────
def slide_problem():
    s = new_slide(OFFWHITE)
    section_head(s, "Utgångsläget", "Föreningskassan behöver fyllas \u2013\nmen sortimentet inspirerar inte")
    cards = [
        ("Engångsköp", "Strumpor, godis, toapapper och rabatthäften säljs en gång. Ingen återkommer av kärlek till varumärket."),
        ("Lågt engagemang", "Det är svårt att känna stolthet över att sälja något man inte själv längtar efter."),
        ("Pressad marginal", "Låg prislapp ger lite per sälj \u2013 det krävs många ordrar för att fylla kassan."),
    ]
    x = Inches(0.7); w = Inches(3.85); gap = Inches(0.18)
    for i, (t, b) in enumerate(cards):
        cx = x + i * (w + gap)
        rrect(s, cx, Inches(2.9), w, Inches(3.1), WHITE, line=SAND_LIGHT)
        dot(s, cx + Inches(0.35), Inches(3.2), Inches(0.32), TERRA)
        txt(s, cx + Inches(0.35), Inches(3.75), w - Inches(0.7), Inches(0.6),
            [[(t, HEAD, 17, INK, True)]], space_after=0)
        txt(s, cx + Inches(0.35), Inches(4.35), w - Inches(0.7), Inches(1.5),
            [[(b, BODY, 12.5, SAND_DARK, False)]], line_spacing=1.3)
    txt(s, Inches(0.7), Inches(6.3), Inches(11.8), Inches(0.7),
        [[("Marknaden är enorm \u2013 Newbody ensamt har samlat in över 1,7 miljarder kr. "
           "Frågan är inte OM man ska sälja, utan VAD.", BODY, 13, FOREST, True)]], line_spacing=1.3)
    footer(s); page_no(s, 2)


# ───────────────────────── 3. LOSNING ─────────────────────────
def slide_solution():
    s = new_slide(INK)
    rect(s, 0, 0, Inches(0.28), SH, FOREST)
    kicker(s, Inches(0.7), Inches(0.7), "Vårt svar", color=OLIVE)
    txt(s, Inches(0.66), Inches(1.15), Inches(11.8), Inches(2.0),
        [[("Sälj något folk använder varje dag \u2013", HEAD, 32, WHITE, True)],
         [("och kommer tillbaka till.", HEAD, 32, OLIVE, True)]],
        space_after=2, line_spacing=1.06)
    pts = [
        ("Premium, inte commodity", "Hudvård och hårvård i butiksklass som kunden faktiskt vill ha."),
        ("Återkommande intäkter", "När flaskan tar slut kommer kunden tillbaka \u2013 försäljning blir flöde, inte engångs."),
        ("Forskningsförankrat", "Certifierade aktiver (NATRUE/COSMOS) och en ärlig berättelse om hudbottens balans."),
        ("Nordiskt och schysst", "Snällt mot håret och naturen \u2013 och varje flaska stöttar er förening."),
    ]
    x = Inches(0.7); y = Inches(3.5); w = Inches(5.7); colgap = Inches(0.5); rowgap = Inches(1.55)
    for i, (t, b) in enumerate(pts):
        cx = x + (i % 2) * (w + colgap)
        cy = y + (i // 2) * rowgap
        dot(s, cx, cy + Inches(0.05), Inches(0.28), OLIVE)
        txt(s, cx + Inches(0.45), cy - Inches(0.05), w - Inches(0.45), Inches(0.5),
            [[(t, HEAD, 16, WHITE, True)]], space_after=0)
        txt(s, cx + Inches(0.45), cy + Inches(0.4), w - Inches(0.45), Inches(1.0),
            [[(b, BODY, 12, SAND_LIGHT, False)]], line_spacing=1.3)
    page_no(s, 3)


def image_slide(kick, title, image, caption=None, n=0, title_color=INK,
                img_max_w=Inches(11.9), img_max_h=Inches(4.3), img_top=Inches(2.4)):
    s = new_slide(OFFWHITE)
    section_head(s, kick, title, title_color=title_color)
    pic_fit(s, image, Inches(0.7), img_top, img_max_w, img_max_h, valign="top")
    if caption:
        txt(s, Inches(0.7), SH - Inches(1.0), Inches(11.9), Inches(0.6),
            [[(caption, BODY, 12, SAND_DARK, False)]], align=PP_ALIGN.CENTER, line_spacing=1.25)
    footer(s); page_no(s, n)
    return s


# 4. Produkter
def slide_products():
    image_slide("Sortimentet", "Tre formuleringar, en filosofi", "product_line.png",
                caption="Rengör mjukt \u00b7 lugnar hudbotten \u00b7 lämnar håret starkare. Sulfatsnålt och biologiskt nedbrytbart.",
                n=4, img_max_h=Inches(4.5))


# 5. Vetenskap (ECS + mikrobiom)
def slide_science_ecs():
    s = new_slide(OFFWHITE)
    section_head(s, "Det som gör oss annorlunda", "Vi jobbar med hudbotten \u2013 inte emot den")
    txt(s, Inches(0.7), Inches(2.15), Inches(11.9), Inches(0.7),
        [[("Håret växer ur en levande hudbotten. Två små system håller den i balans \u2013 och Roots är gjort för att stötta dem.",
           BODY, 14, SAND_DARK, False)]], line_spacing=1.3)
    pic_fit(s, "ecs_microbiome.png", Inches(0.7), Inches(2.95), Inches(11.9), Inches(3.6), valign="top")
    footer(s); page_no(s, 5)


# 6. Balans-diagram
def slide_balance():
    s = new_slide(OFFWHITE)
    section_head(s, "Tre pelare", "Lugn, fukt och skydd \u2013 där håret börjar")
    pic_fit(s, "science_balance.png", Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.5), valign="top")
    footer(s); page_no(s, 6)


# 7. Ingredienser
def slide_ingredients():
    image_slide("Nyckelingredienser", "Få ingredienser \u2013 de som finns gör jobbet",
                "ingredient_cards.png",
                caption="SyriCalm och MultiMoist är certifierade enligt NATRUE och COSMOS \u2013 oberoende standarder för naturlig kosmetik.",
                n=7, img_max_h=Inches(4.2))


# 8. Positionering
def slide_positioning():
    s = new_slide(OFFWHITE)
    section_head(s, "Varför det säljer bättre", "En öppen flank på en stor marknad")
    pic_fit(s, "positioning_map.png", Inches(0.6), Inches(2.2), Inches(6.6), Inches(4.7), align="left", valign="top")
    pts = [
        ("Högre ordervärde", "En riktigt bra produkt ger större ordrar \u2013 mer i kassan per sälj."),
        ("Återkommande köp", "Hudvård tar slut. Kunden kommer tillbaka \u2013 utrymme för prenumeration."),
        ("Lätt att sälja", "Det är roligare att sälja något man är stolt över. Stolthet smittar."),
    ]
    x = Inches(7.5); y = Inches(2.5)
    for i, (t, b) in enumerate(pts):
        cy = y + i * Inches(1.4)
        dot(s, x, cy + Inches(0.04), Inches(0.26), FOREST)
        txt(s, x + Inches(0.42), cy - Inches(0.06), Inches(5.0), Inches(0.5),
            [[(t, HEAD, 17, INK, True)]], space_after=0)
        txt(s, x + Inches(0.42), cy + Inches(0.4), Inches(5.0), Inches(0.9),
            [[(b, BODY, 12.5, SAND_DARK, False)]], line_spacing=1.3)
    footer(s); page_no(s, 8)


# 9. Varde till kassan
def slide_value():
    image_slide("Vad det ger föreningen", "Premium ger mer i kassan över tid",
                "value_compare.png",
                caption="Illustrativt: även med rimlig marginal ger högre ordervärde + återköp mer än ett billigt engångsköp.",
                n=9, img_max_h=Inches(4.0))


# 10. Sa funkar det
def slide_howitworks():
    s = new_slide(OFFWHITE)
    section_head(s, "Så går det till", "Digitalt, enkelt och utan risk")
    pic_fit(s, "sales_journey.png", Inches(0.7), Inches(2.5), Inches(11.9), Inches(2.4), valign="top")
    chips = ["Ingen startavgift", "Personlig länk + QR", "Klarna & Swish", "Realtid i panelen", "Faktura efter leverans"]
    x = Inches(0.7); y = Inches(5.6)
    for c in chips:
        w = Inches(0.18 + 0.115 * len(c))
        rrect(s, x, y, w, Inches(0.55), SAND_100)
        txt(s, x, y + Inches(0.04), w, Inches(0.45),
            [[(c, BODY, 11.5, INK, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        x = x + w + Inches(0.2)
    footer(s); page_no(s, 10)


# 11. Nordiskt och schysst
def slide_values():
    s = new_slide(INK)
    rect(s, 0, 0, Inches(0.28), SH, FOREST)
    kicker(s, Inches(0.7), Inches(0.9), "Vårt varför", color=OLIVE)
    txt(s, Inches(0.66), Inches(1.4), Inches(11.8), Inches(2.2),
        [[("Vi vill bevisa att man kan", HEAD, 30, WHITE, True)],
         [("göra gott och sälja riktigt bra saker.", HEAD, 30, OLIVE, True)]],
        space_after=2, line_spacing=1.08)
    txt(s, Inches(0.7), Inches(3.5), Inches(11.5), Inches(2.2),
        [[("Varje flaska som säljs ger pengar till en förening i Sverige \u2013 och en bättre dag för "
           "någons hår. Inga genvägar på kvaliteten, inga överdrifter i löftena. Bara nordisk omtanke, "
           "från hudbotten och ut.", BODY, 17, SAND_LIGHT, False)]], line_spacing=1.4)
    page_no(s, 11)


# 12. CTA
def slide_cta():
    s = new_slide(OFFWHITE)
    rect(s, 0, 0, SW, Inches(0.28), FOREST)
    pic_fit(s, "brand/roots-symbol-black.png", Inches(6.0), Inches(1.1),
            Inches(1.3), Inches(1.3), from_web=True)
    txt(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.4),
        [[("Ska vi fylla er kassa \u2013", HEAD, 34, INK, True)],
         [("och ge era medlemmar bättre hår?", HEAD, 34, FOREST, True)]],
        align=PP_ALIGN.CENTER, space_after=2, line_spacing=1.06)
    txt(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.0),
        [[("Vi hjälper er igång på under fem minuter. Ingen startavgift, ingen bindningstid \u2013 "
           "ni betalar först när produkterna är sålda och levererade.", BODY, 15, SAND_DARK, False)]],
        align=PP_ALIGN.CENTER, line_spacing=1.35)
    rrect(s, Inches(5.16), Inches(5.9), Inches(3.0), Inches(0.7), FOREST)
    txt(s, Inches(5.16), Inches(5.99), Inches(3.0), Inches(0.55),
        [[("Starta er försäljning", BODY, 14, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    txt(s, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.4),
        [[("roots.se   \u00b7   hej@roots.se", BODY, 12, SAND_DARK, False)]],
        align=PP_ALIGN.CENTER, space_after=0)


def build():
    slide_cover()
    slide_problem()
    slide_solution()
    slide_products()
    slide_science_ecs()
    slide_balance()
    slide_ingredients()
    slide_positioning()
    slide_value()
    slide_howitworks()
    slide_values()
    slide_cta()
    prs.save(OUT)
    print("Saved", os.path.relpath(OUT, ROOT), "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")


build()
